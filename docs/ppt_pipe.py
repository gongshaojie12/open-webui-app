import os
import re
import time
import asyncio
import threading  # 【新增】OBS 客户端初始化锁
from concurrent.futures import ThreadPoolExecutor  # 【新增】OBS 上传专用线程池
import httpx
import base64
import hashlib
import logging
import json
import io
import uuid
import aiofiles
import mimetypes
from PIL import Image
from google import genai
from google.genai import types
from google.genai.errors import ClientError, ServerError, APIError
from typing import (
    List,
    Union,
    Optional,
    Dict,
    Any,
    Tuple,
    AsyncIterator,
    Callable,
    Literal,
)
from pydantic_core import core_schema
from pydantic import BaseModel, Field, GetCoreSchemaHandler
from cryptography.fernet import Fernet, InvalidToken

# from open_webui.env import SRC_LOG_LEVELS
from fastapi import Request, UploadFile, BackgroundTasks
from open_webui.routers.files import upload_file_handler
from open_webui.models.users import UserModel, Users

try:
    from open_webui.socket.main import sio as _owui_sio
except ImportError:
    _owui_sio = None
from starlette.datastructures import Headers
from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
import aiohttp
from obs import ObsClient, PutObjectHeader  # 【新增】华为 OBS SDK

# 【新增】OBS 上传专用线程池（延迟初始化，由 Valves.OBS_UPLOAD_WORKERS 控制大小）
_obs_executor: Optional[ThreadPoolExecutor] = None
_obs_executor_size: int = 0
_obs_executor_init_lock = threading.Lock()


class Pipe:
    """
    Pipeline for interacting with Google Gemini models.
    """

    # Configuration valves for the pipeline
    class Valves(BaseModel):
        APP_TYPE: Literal[
            "chatflow",
            "workflow",
        ] = Field(
            default="chatflow",
            description="Dify API calling method.",
        )
        BASE_URL: str = Field(
            default=os.getenv("BASE_URL", "http://192.168.4.114"),
            description="Dify Base URL for the PPT AI API.",
        )
        OPEN_WEBUI_BASE_URL: str = Field(
            default=os.getenv("OPEN_WEBUI_BASE_URL", "http://192.168.4.114:8080"),
            description="This URL is used to access the content of uploaded documents..",
        )
        API_KEY: str = Field(
            default=os.getenv("API_KEY", ""),
            description="API key for PPT Generative AI.",
        )
        WORKFLOW_ID: str = Field(
            default=os.getenv("WORKFLOW_ID", ""),
            description="API key for PPT Generative AI.",
        )
        # PPT_TYPE: Literal["ppt", "pdf"] = Field(
        #     default="pdf",
        #     description="File type to generate.",
        # )
        IMAGE_ENABLE_OPTIMIZATION: bool = Field(
            default=os.getenv("IMAGE_ENABLE_OPTIMIZATION", "true").lower() == "true",
            description="Enable intelligent image optimization for API compatibility",
        )
        IMAGE_MAX_SIZE_MB: float = Field(
            default=float(os.getenv("IMAGE_MAX_SIZE_MB", "15.0")),
            description="Maximum image size in MB before compression is applied",
        )
        IMAGE_MAX_DIMENSION: int = Field(
            default=int(os.getenv("IMAGE_MAX_DIMENSION", "2048")),
            description="Maximum width or height in pixels before resizing",
        )
        IMAGE_COMPRESSION_QUALITY: int = Field(
            default=int(os.getenv("IMAGE_COMPRESSION_QUALITY", "85")),
            description="JPEG compression quality (1-100, higher = better quality but larger size)",
        )

        IMAGE_PNG_COMPRESSION_THRESHOLD_MB: float = Field(
            default=float(os.getenv("IMAGE_PNG_THRESHOLD_MB", "0.5")),
            description="PNG files above this size (MB) will be converted to JPEG for better compression",
        )

        # 【新增】─── OBS 配置（参考 seedance.py） ───
        OBS_SERVER: str = Field(
            default="https://obs.cn-north-4.myhuaweicloud.com",
            description="华为 OBS Endpoint（如 https://obs.cn-north-4.myhuaweicloud.com）",
        )
        OBS_BUCKET_NAME: str = Field(
            default="",
            description="OBS Bucket 名称",
        )
        OBS_ACCESS_KEY: str = Field(
            default="",
            description="OBS Access Key ID（AK）",
        )
        OBS_SECRET_KEY: str = Field(
            default="",
            description="OBS Secret Access Key（SK）",
        )
        OBS_UPLOAD_WORKERS: int = Field(
            default=16,
            description="OBS 上传线程池大小（并发上传数上限）",
            ge=1,
            le=128,
        )
        OBS_UPLOAD_MAX_RETRIES: int = Field(
            default=3,
            description="OBS 上传失败最大重试次数",
            ge=1,
            le=10,
        )

    # ---------------- Internal Helpers ---------------- #
    async def _gather_history_images(
        self,
        messages: List[Dict[str, Any]],
        last_user_msg: Dict[str, Any],
        optimization_stats: List[Dict[str, Any]],
    ) -> List[Dict[str, Any]]:
        history_images: List[Dict[str, Any]] = []
        for msg in messages:
            if msg is last_user_msg:
                continue
            if msg.get("role") not in {"user", "assistant"}:
                continue
            _p, parts = await self._extract_images_from_message(
                msg, stats_list=optimization_stats
            )
            if parts:
                history_images.extend(parts)
        return history_images

    def _deduplicate_images(self, images: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        # if not self.valves.IMAGE_DEDUP_HISTORY:
        #     return images
        seen: set[str] = set()
        result: List[Dict[str, Any]] = []
        for part in images:
            try:
                data = part["inline_data"]["data"]
                # Hash full base64 payload for stronger dedup reliability
                h = hashlib.sha256(data.encode()).hexdigest()
                if h in seen:
                    continue
                seen.add(h)
            except Exception:
                pass
            result.append(part)
        return result

    def _apply_order_and_limit(
        self,
        history: List[Dict[str, Any]],
        current: List[Dict[str, Any]],
    ) -> Tuple[List[Dict[str, Any]], List[bool]]:
        """Combine history & current image parts honoring order & global limit.

        Returns:
            (combined_parts, reused_flags) where reused_flags[i] == True indicates
            the image originated from history, False if from current message.
        """
        history_first = self.valves.IMAGE_HISTORY_FIRST
        limit = max(1, self.valves.IMAGE_HISTORY_MAX_REFERENCES)
        combined: List[Dict[str, Any]] = []
        reused_flags: List[bool] = []

        def append(parts: List[Dict[str, Any]], reused: bool):
            for p in parts:
                if len(combined) >= limit:
                    break
                combined.append(p)
                reused_flags.append(reused)

        if history_first:
            append(history, True)
            append(current, False)
        else:
            append(current, False)
            append(history, True)
        return combined, reused_flags

    async def _emit_image_stats(
        self,
        ordered_stats: List[Dict[str, Any]],
        reused_flags: List[bool],
        total_limit: int,
        __event_emitter__: Callable,
    ) -> None:
        """Emit per-image optimization stats aligned with final combined order.

        ordered_stats: stats list in the exact order images will be sent (same length as combined image list)
        reused_flags: parallel list indicating whether image originated from history
        """
        if not ordered_stats:
            return
        for idx, stat in enumerate(ordered_stats, start=1):
            reused = reused_flags[idx - 1] if idx - 1 < len(reused_flags) else False
            stat_copy = dict(stat) if stat else {}
            stat_copy.update({"index": idx, "reused": reused})
            if stat and stat.get("original_size_mb") is not None:
                desc = f"Image {idx}: {stat['original_size_mb']:.2f}MB -> {stat['final_size_mb']:.2f}MB"
                if stat.get("quality") is not None:
                    desc += f" (Q{stat['quality']})"
            else:
                desc = f"Image {idx}: (no metrics)"
            reasons = stat.get("reasons") if stat else None
            if reasons:
                desc += " | " + ", ".join(reasons[:3])
            await __event_emitter__(
                {
                    "type": "status",
                    "data": {
                        "action": "image_optimization",
                        "description": desc,
                        "index": idx,
                        "done": False,
                        "details": stat_copy,
                    },
                }
            )
        await __event_emitter__(
            {
                "type": "status",
                "data": {
                    "action": "image_optimization",
                    "description": f"{len(ordered_stats)} image(s) processed (limit {total_limit}).",
                    "done": True,
                },
            }
        )

    async def _build_image_generation_contents(
        self,
        messages: List[Dict[str, Any]],
        __event_emitter__: Callable,
    ) -> Tuple[List[Dict[str, Any]], Optional[str]]:
        """Construct the contents payload for image-capable models.

        Returns tuple (contents, system_instruction) where system_instruction is extracted from system messages.
        """
        # Extract system instruction first
        system_instruction = next(
            (msg["content"] for msg in messages if msg.get("role") == "system"),
            None,
        )

        last_user_msg = next(
            (m for m in reversed(messages) if m.get("role") == "user"), None
        )
        if not last_user_msg:
            raise ValueError("No user message found")

        optimization_stats: List[Dict[str, Any]] = []
        history_images = await self._gather_history_images(
            messages, last_user_msg, optimization_stats
        )
        prompt, current_images = await self._extract_images_from_message(
            last_user_msg, stats_list=optimization_stats
        )

        # Deduplicate
        history_images = self._deduplicate_images(history_images)
        current_images = self._deduplicate_images(current_images)

        combined, reused_flags = self._apply_order_and_limit(
            history_images, current_images
        )

        if not prompt and not combined:
            raise ValueError("No prompt or images provided")
        if not prompt and combined:
            prompt = "Analyze and describe the provided images."

        # Build ordered stats aligned with combined list
        ordered_stats: List[Dict[str, Any]] = []
        if optimization_stats:
            # Build map from final_hash -> stat (first wins)
            hash_map: Dict[str, Dict[str, Any]] = {}
            for s in optimization_stats:
                fh = s.get("final_hash")
                if fh and fh not in hash_map:
                    hash_map[fh] = s
            for part in combined:
                try:
                    fh = hashlib.sha256(
                        part["inline_data"]["data"].encode()
                    ).hexdigest()
                    ordered_stats.append(hash_map.get(fh) or {})
                except Exception:
                    ordered_stats.append({})
        # Emit stats AFTER final ordering so labels match
        await self._emit_image_stats(
            ordered_stats,
            reused_flags,
            self.valves.IMAGE_HISTORY_MAX_REFERENCES,
            __event_emitter__,
        )

        # Emit mapping
        if combined:
            mapping = [
                {
                    "index": i + 1,
                    "label": (
                        f"Image {i + 1}" if self.valves.IMAGE_ADD_LABELS else str(i + 1)
                    ),
                    "reused": reused_flags[i],
                    "origin": "history" if reused_flags[i] else "current",
                }
                for i in range(len(combined))
            ]
            await __event_emitter__(
                {
                    "type": "status",
                    "data": {
                        "action": "image_reference_map",
                        "description": f"{len(combined)} image(s) included (limit {self.valves.IMAGE_HISTORY_MAX_REFERENCES}).",
                        "images": mapping,
                        "done": True,
                    },
                }
            )

        # Build parts
        parts: List[Dict[str, Any]] = []

        # For image generation models, prepend system instruction to the prompt
        # since system_instruction parameter may not be supported
        final_prompt = prompt
        if system_instruction and prompt:
            final_prompt = f"{system_instruction}\n\n{prompt}"
            self.log.debug(
                f"Prepended system instruction to prompt for image generation. "
                f"System instruction length: {len(system_instruction)}, "
                f"Original prompt length: {len(prompt)}, "
                f"Final prompt length: {len(final_prompt)}"
            )
        elif system_instruction and not prompt:
            final_prompt = system_instruction
            self.log.debug(
                f"Using system instruction as prompt for image generation "
                f"(length: {len(system_instruction)})"
            )

        if final_prompt:
            parts.append({"text": final_prompt})
        if self.valves.IMAGE_ADD_LABELS:
            for idx, part in enumerate(combined, start=1):
                parts.append({"text": f"[Image {idx}]"})
                parts.append(part)
        else:
            parts.extend(combined)

        self.log.debug(
            f"Image-capable payload: history={len(history_images)} current={len(current_images)} used={len(combined)} limit={self.valves.IMAGE_HISTORY_MAX_REFERENCES} history_first={self.valves.IMAGE_HISTORY_FIRST} prompt_len={len(final_prompt)}"
        )
        # Return None for system_instruction since we've incorporated it into the prompt
        return [{"role": "user", "parts": parts}], None

    def __init__(self):
        """Initializes the Pipe instance and configures the genai library."""
        self.valves = self.Valves()
        self.name: str = "PPT:"

        # Setup logging
        self.log = logging.getLogger("google_ai.pipe")
        self.log.setLevel("INFO")

        # 【新增】OBS 客户端（延迟初始化）
        self._obs_client = None
        self._obs_lock = threading.Lock()

        # # Model cache
        # self._model_cache: Optional[List[Dict[str, str]]] = None
        # self._model_cache_time: float = 0

    # 【新增】─── OBS 操作（参考 seedance.py） ───

    def _get_obs_client(self) -> ObsClient:
        if self._obs_client is None:
            with self._obs_lock:
                if self._obs_client is None:
                    self._obs_client = ObsClient(
                        access_key_id=self.valves.OBS_ACCESS_KEY,
                        secret_access_key=self.valves.OBS_SECRET_KEY,
                        server=self.valves.OBS_SERVER,
                    )
        return self._obs_client

    @staticmethod
    def _get_obs_executor(workers: int) -> ThreadPoolExecutor:
        """按 Valves 配置动态创建/扩容线程池（进程内共享）。"""
        global _obs_executor, _obs_executor_size
        if _obs_executor is None or _obs_executor_size != workers:
            with _obs_executor_init_lock:
                if _obs_executor is None or _obs_executor_size != workers:
                    _obs_executor = ThreadPoolExecutor(
                        max_workers=workers, thread_name_prefix="obs-upload"
                    )
                    _obs_executor_size = workers
        return _obs_executor

    def _upload_to_obs(
        self,
        data: bytes,
        object_key: str,
        content_type: str = "application/octet-stream",
    ) -> str:
        """同步上传文件到华为 OBS，返回公网 URL。Bucket 须设为公共读。"""
        client = self._get_obs_client()
        bucket = self.valves.OBS_BUCKET_NAME
        max_retries = self.valves.OBS_UPLOAD_MAX_RETRIES
        last_error = None

        headers = PutObjectHeader()
        headers.contentType = content_type

        for attempt in range(1, max_retries + 1):
            try:
                resp = client.putContent(
                    bucketName=bucket,
                    objectKey=object_key,
                    content=io.BytesIO(data),
                    headers=headers,
                )
                if resp.status < 300:
                    server = self.valves.OBS_SERVER.replace("https://", "").replace(
                        "http://", ""
                    )
                    url = f"https://{bucket}.{server}/{object_key}"
                    return url
                else:
                    raise Exception(
                        f"OBS 上传失败 (status={resp.status}): {resp.errorCode} - {resp.errorMessage}"
                    )
            except Exception as e:
                last_error = e
                self.log.warning(f"OBS 上传第 {attempt}/{max_retries} 次失败: {e}")
                if attempt < max_retries:
                    time.sleep(1 * attempt)

        raise Exception(f"OBS 上传失败（重试 {max_retries} 次）: {last_error}")

    async def _upload_to_obs_async(
        self,
        data: bytes,
        object_key: str,
        content_type: str = "application/octet-stream",
    ) -> str:
        """异步包装 OBS 上传（OBS SDK 是同步的，用可配置大小的专用线程池执行）。"""
        loop = asyncio.get_running_loop()
        executor = self._get_obs_executor(self.valves.OBS_UPLOAD_WORKERS)
        return await loop.run_in_executor(
            executor, self._upload_to_obs, data, object_key, content_type
        )

    # ─── 原有方法 ───

    def _get_client(self) -> genai.Client:
        """
        Validates API credentials and returns a genai.Client instance.
        """
        self._validate_api_key()

        if self.valves.USE_VERTEX_AI:
            self.log.debug(
                f"Initializing Vertex AI client (Project: {self.valves.VERTEX_PROJECT}, Location: {self.valves.VERTEX_LOCATION})"
            )
            return genai.Client(
                vertexai=True,
                project=self.valves.VERTEX_PROJECT,
                location=self.valves.VERTEX_LOCATION,
            )
        else:
            self.log.debug("Initializing Google Generative AI client with API Key")
            options = types.HttpOptions(
                api_version=self.valves.API_VERSION, base_url=self.valves.BASE_URL
            )
            return genai.Client(
                api_key=self.valves.API_KEY,
                http_options=options,
            )

    def _validate_api_key(self) -> None:
        """
        Validates that the necessary Google API credentials are set.

        Raises:
            ValueError: If the required credentials are not set.
        """
        if self.valves.USE_VERTEX_AI:
            if not self.valves.VERTEX_PROJECT:
                self.log.error("USE_VERTEX_AI is true, but VERTEX_PROJECT is not set.")
                raise ValueError(
                    "VERTEX_PROJECT is not set. Please provide the Google Cloud project ID."
                )
            # For Vertex AI, location has a default, so project is the main thing to check.
            # Actual authentication will be handled by ADC or environment.
            self.log.debug(
                "Using Vertex AI. Ensure ADC or service account is configured."
            )
        else:
            if not self.valves.API_KEY:
                self.log.error("API_KEY is not set (and not using Vertex AI).")
                raise ValueError(
                    "API_KEY is not set. Please provide the API key in the environment variables or valves."
                )
            self.log.debug("Using Google Generative AI API with API Key.")

    def strip_prefix(self, model_name: str) -> str:
        """
        Extract the model identifier using regex, handling various naming conventions.
        e.g., "google_gemini_pipeline.gemini-2.5-flash-preview-04-17" -> "gemini-2.5-flash-preview-04-17"
        e.g., "models/gemini-1.5-flash-001" -> "gemini-1.5-flash-001"
        e.g., "publishers/google/models/gemini-1.5-pro" -> "gemini-1.5-pro"
        """
        # Use regex to remove everything up to and including the last '/' or the first '.'
        stripped = re.sub(r"^(?:.*/|[^.]*\.)", "", model_name)
        return stripped

    def get_google_models(self, force_refresh: bool = False) -> List[Dict[str, str]]:
        """
        Retrieve available Google models suitable for content generation.
        Uses caching to reduce API calls.

        Args:
            force_refresh: Whether to force refreshing the model cache

        Returns:
            List of dictionaries containing model id and name.
        """
        # Check cache first
        current_time = time.time()
        if (
            not force_refresh
            and self._model_cache is not None
            and (current_time - self._model_cache_time) < self.valves.MODEL_CACHE_TTL
        ):
            self.log.debug("Using cached model list")
            return self._model_cache

        try:
            client = self._get_client()
            models = client.models.list()
            model_map: Dict[str, Dict[str, Any]] = {}
            for model in models:
                model_id = self.strip_prefix(getattr(model, "name", ""))
                if not model_id.startswith("gemini-"):
                    continue
                display = getattr(model, "display_name", None)
                supports_image_generation = self._check_image_generation_support(
                    model_id
                )
                model_name = display or model_id
                if supports_image_generation:
                    model_name += " ??"
                model_map[model_id] = {
                    "id": model_id,
                    "name": model_name,
                    "image_generation": supports_image_generation,
                }

            if self.valves.USE_VERTEX_AI:
                if not any(k.startswith("gemini-3") for k in model_map.keys()):
                    fallback_locations = [
                        "us-east5",
                        "us-central1",
                        "europe-west1",
                        "asia-southeast1",
                    ]
                    for loc in fallback_locations:
                        try:
                            alt_client = genai.Client(
                                vertexai=True,
                                project=self.valves.VERTEX_PROJECT,
                                location=loc,
                            )
                            alt_models = alt_client.models.list()
                            for m in alt_models:
                                mid = self.strip_prefix(getattr(m, "name", ""))
                                if not mid.startswith("gemini-"):
                                    continue
                                disp = getattr(m, "display_name", None)
                                img = self._check_image_generation_support(mid)
                                name = disp or mid
                                if img:
                                    name += " ??"
                                if mid not in model_map:
                                    model_map[mid] = {
                                        "id": mid,
                                        "name": name,
                                        "image_generation": img,
                                    }
                        except Exception:
                            pass

            builtin = [
                m.strip()
                for m in (self.valves.BUILTIN_MODELS or "").split(",")
                if m.strip()
            ]
            for mid in builtin:
                self.log.debug(
                    f"mid================================================{mid}"
                )
                if not mid.startswith("gemini-"):
                    continue
                if "image" not in mid:  # 只返回图片模型
                    continue
                if mid not in model_map:
                    img = self._check_image_generation_support(mid)
                    name = mid.replace("gemini-", "Gemini ").replace("-", " ").title()
                    if img:
                        name += " ??"
                    model_map[mid] = {"id": mid, "name": name, "image_generation": img}

            self._model_cache = list(model_map.values())
            self._model_cache_time = current_time
            return self._model_cache

        except Exception as e:
            self.log.exception(f"Could not fetch models from Google: {str(e)}")
            # Return a specific error entry for the UI
            return [{"id": "error", "name": f"Could not fetch models: {str(e)}"}]

    def _check_image_generation_support(self, model_id: str) -> bool:
        """
        Check if a model supports image generation.

        Args:
            model_id: The model ID to check

        Returns:
            True if the model supports image generation, False otherwise
        """
        # Known image generation models
        image_generation_models = [
            "gemini-2.5-flash-image-preview",
            "gemini-2.5-flash-image",
        ]

        # Check for exact matches or pattern matches
        for pattern in image_generation_models:
            if model_id == pattern or pattern in model_id:
                return True

        # Additional pattern checking for future models
        if "image" in model_id.lower() and (
            "generation" in model_id.lower() or "preview" in model_id.lower()
        ):
            return True

        return False

    def _check_thinking_support(self, model_id: str) -> bool:
        """
        Check if a model supports the thinking feature.

        Args:
            model_id: The model ID to check

        Returns:
            True if the model supports thinking, False otherwise
        """
        # Models that do NOT support thinking
        non_thinking_models = [
            "gemini-2.5-flash-image-preview",
            "gemini-2.5-flash-image",
        ]

        # Check for exact matches
        for pattern in non_thinking_models:
            if model_id == pattern or pattern in model_id:
                return False

        # Additional pattern checking - image generation models typically don't support thinking
        if "image" in model_id.lower() and (
            "generation" in model_id.lower() or "preview" in model_id.lower()
        ):
            return False

        # By default, assume models support thinking
        return True

    def _check_thinking_level_support(self, model_id: str) -> bool:
        """
        Check if a model supports the thinking_level parameter.

        Gemini 3 models support thinking_level and should NOT use thinking_budget.
        Other models (like Gemini 2.5) use thinking_budget instead.

        Args:
            model_id: The model ID to check

        Returns:
            True if the model supports thinking_level, False otherwise
        """
        # Gemini 3 models support thinking_level (not thinking_budget)
        gemini_3_patterns = [
            "gemini-3-",
        ]

        model_lower = model_id.lower()
        for pattern in gemini_3_patterns:
            if pattern in model_lower:
                return True

        return False

    def _validate_thinking_level(self, level: str) -> Optional[str]:
        """
        Validate and normalize the thinking level value.

        Args:
            level: The thinking level string to validate

        Returns:
            Normalized level string ('low', 'high') or None if invalid/empty
        """
        if not level:
            return None

        normalized = level.strip().lower()
        valid_levels = ["low", "high"]

        if normalized in valid_levels:
            return normalized

        self.log.warning(
            f"Invalid thinking level '{level}'. Valid values are: {', '.join(valid_levels)}. "
            "Falling back to model default."
        )
        return None

    def _validate_thinking_budget(self, budget: int) -> int:
        """
        Validate and normalize the thinking budget value.

        Args:
            budget: The thinking budget integer to validate

        Returns:
            Validated budget: -1 for dynamic, 0 to disable, or 1-32768 for fixed limit
        """
        # -1 means dynamic thinking (let the model decide)
        if budget == -1:
            return -1

        # 0 means disable thinking
        if budget == 0:
            return 0

        # Validate positive range (1-32768)
        if budget > 0:
            if budget > 32768:
                self.log.warning(
                    f"Thinking budget {budget} exceeds maximum of 32768. Clamping to 32768."
                )
                return 32768
            return budget

        # Negative values (except -1) are invalid, treat as -1 (dynamic)
        self.log.warning(
            f"Invalid thinking budget {budget}. Only -1 (dynamic), 0 (disabled), or 1-32768 are valid. "
            "Falling back to dynamic thinking."
        )
        return -1

    def pipes(self) -> List[Dict[str, str]]:
        return [{"id": "ai-ppt", "name": "PPT Generate AI"}]

    def _prepare_model_id(self, model_id: str) -> str:
        """
        Prepare and validate the model ID for use with the API.

        Args:
            model_id: The original model ID from the user

        Returns:
            Properly formatted model ID

        Raises:
            ValueError: If the model ID is invalid or unsupported
        """
        original_model_id = model_id
        model_id = self.strip_prefix(model_id)

        # If the model ID doesn't look like a Gemini model, try to find it by name
        if not model_id.startswith("gemini-"):
            models_list = self.get_google_models()
            found_model = next(
                (m["id"] for m in models_list if m["name"] == original_model_id), None
            )
            if found_model and found_model.startswith("gemini-"):
                model_id = found_model
                self.log.debug(
                    f"Mapped model name '{original_model_id}' to model ID '{model_id}'"
                )
            else:
                # If we still don't have a valid ID, raise an error
                if not model_id.startswith("gemini-"):
                    self.log.error(
                        f"Invalid or unsupported model ID: '{original_model_id}'"
                    )
                    raise ValueError(
                        f"Invalid or unsupported Google model ID or name: '{original_model_id}'"
                    )

        return model_id

    def _prepare_content(
        self, messages: List[Dict[str, Any]]
    ) -> Tuple[List[Dict[str, Any]], Optional[str]]:
        """
        Prepare messages content for the API and extract system message if present.

        Args:
            messages: List of message objects from the request

        Returns:
            Tuple of (prepared content list, system message string or None)
        """
        # Extract system message
        system_message = next(
            (msg["content"] for msg in messages if msg.get("role") == "system"),
            None,
        )

        # Prepare contents for the API
        contents = []
        for message in messages:
            role = message.get("role")
            if role == "system":
                continue  # Skip system messages, handled separately

            content = message.get("content", "")
            parts = []

            # Handle different content types
            if isinstance(content, list):  # Multimodal content
                parts.extend(self._process_multimodal_content(content))
            elif isinstance(content, str):  # Plain text content
                parts.append({"text": content})
            else:
                self.log.warning(f"Unsupported message content type: {type(content)}")
                continue  # Skip unsupported content

            # Map roles: 'assistant' -> 'model', 'user' -> 'user'
            api_role = "model" if role == "assistant" else "user"
            if parts:  # Only add if there are parts
                contents.append({"role": api_role, "parts": parts})

        return contents, system_message

    def _process_multimodal_content(
        self, content_list: List[Dict[str, Any]]
    ) -> List[Dict[str, Any]]:
        """
        Process multimodal content (text and images).

        Args:
            content_list: List of content items

        Returns:
            List of processed parts for the Gemini API
        """
        parts = []

        for item in content_list:
            if item.get("type") == "text":
                parts.append({"text": item.get("text", "")})
            elif item.get("type") == "image_url":
                image_url = item.get("image_url", {}).get("url", "")

                if image_url.startswith("data:image"):
                    # Handle base64 encoded image data with optimization
                    try:
                        # Optimize the image before processing
                        optimized_image = self._optimize_image_for_api(image_url)
                        header, encoded = optimized_image.split(",", 1)
                        mime_type = header.split(":")[1].split(";")[0]

                        # Basic validation for image types
                        if mime_type not in [
                            "image/jpeg",
                            "image/png",
                            "image/webp",
                            "image/heic",
                            "image/heif",
                        ]:
                            self.log.warning(
                                f"Unsupported image mime type: {mime_type}"
                            )
                            parts.append(
                                {"text": f"[Image type {mime_type} not supported]"}
                            )
                            continue

                        # Check if the encoded data is too large
                        if len(encoded) > 15 * 1024 * 1024:  # 15MB limit for base64
                            self.log.warning(
                                f"Image data too large: {len(encoded)} characters"
                            )
                            parts.append(
                                {
                                    "text": "[Image too large for processing - please use a smaller image]"
                                }
                            )
                            continue

                        parts.append(
                            {
                                "inline_data": {
                                    "mime_type": mime_type,
                                    "data": encoded,
                                }
                            }
                        )
                    except Exception as img_ex:
                        self.log.exception(f"Could not parse image data URL: {img_ex}")
                        parts.append({"text": "[Image data could not be processed]"})
                else:
                    # Gemini API doesn't directly support image URLs
                    self.log.warning(f"Direct image URLs not supported: {image_url}")
                    parts.append({"text": f"[Image URL not processed: {image_url}]"})

        return parts

    # _find_image removed (was single-image oriented and is superseded by multi-image logic)

    async def _extract_images_from_message(
        self,
        message: Dict[str, Any],
        *,
        stats_list: Optional[List[Dict[str, Any]]] = None,
    ) -> Tuple[str, List[Dict[str, Any]]]:
        """Extract prompt text and ALL images from a single user message.

        This replaces the previous single-image _find_image logic for image-capable
        models so that multi-image prompts are respected.

        Returns:
            (prompt_text, image_parts)
                prompt_text: concatenated text content (may be empty)
                image_parts: list of {"inline_data": {mime_type, data}} dicts
        """
        content = message.get("content", "")
        text_segments: List[str] = []
        image_parts: List[Dict[str, Any]] = []

        # Helper to process a data URL or fetched file and append inline_data
        def _add_image(data_url: str):
            try:
                optimized = self._optimize_image_for_api(data_url, stats_list)
                header, b64 = optimized.split(",", 1)
                mime = header.split(":", 1)[1].split(";", 1)[0]
                image_parts.append({"inline_data": {"mime_type": mime, "data": b64}})
            except Exception as e:  # pragma: no cover - defensive
                self.log.warning(f"Skipping image (parse failure): {e}")

        # Regex to extract markdown image references
        md_pattern = re.compile(
            r"!\[[^\]]*\]\((data:image[^)]+|/files/[^)]+|/api/v1/files/[^)]+)\)"
        )

        # Structured multimodal array
        if isinstance(content, list):
            for item in content:
                if item.get("type") == "text":
                    txt = item.get("text", "")
                    text_segments.append(txt)
                    # Also parse any markdown images embedded in the text
                    for match in md_pattern.finditer(txt):
                        url = match.group(1)
                        if url.startswith("data:"):
                            _add_image(url)
                        else:
                            b64 = await self._fetch_file_as_base64(url)
                            if b64:
                                _add_image(b64)
                elif item.get("type") == "image_url":
                    url = item.get("image_url", {}).get("url", "")
                    if url.startswith("data:"):
                        _add_image(url)
                    elif "/files/" in url or "/api/v1/files/" in url:
                        b64 = await self._fetch_file_as_base64(url)
                        if b64:
                            _add_image(b64)
        # Plain string message (may include markdown images)
        elif isinstance(content, str):
            text_segments.append(content)
            for match in md_pattern.finditer(content):
                url = match.group(1)
                if url.startswith("data:"):
                    _add_image(url)
                else:
                    b64 = await self._fetch_file_as_base64(url)
                    if b64:
                        _add_image(b64)
        else:
            self.log.debug(
                f"Unsupported content type for image extraction: {type(content)}"
            )

        prompt_text = " ".join(s.strip() for s in text_segments if s.strip())
        return prompt_text, image_parts

    def _optimize_image_for_api(
        self, image_data: str, stats_list: Optional[List[Dict[str, Any]]] = None
    ) -> str:
        """
        Optimize image data for Gemini API using configurable parameters.

        Returns:
            Optimized base64 data URL
        """
        # Check if optimization is enabled
        if not self.valves.IMAGE_ENABLE_OPTIMIZATION:
            self.log.debug("Image optimization disabled via configuration")
            return image_data

        max_size_mb = self.valves.IMAGE_MAX_SIZE_MB
        max_dimension = self.valves.IMAGE_MAX_DIMENSION
        base_quality = self.valves.IMAGE_COMPRESSION_QUALITY
        png_threshold = self.valves.IMAGE_PNG_COMPRESSION_THRESHOLD_MB

        self.log.debug(
            f"Image optimization config: max_size={max_size_mb}MB, max_dim={max_dimension}px, quality={base_quality}, png_threshold={png_threshold}MB"
        )
        try:
            # Parse the data URL
            if image_data.startswith("data:"):
                header, encoded = image_data.split(",", 1)
                mime_type = header.split(":")[1].split(";")[0]
            else:
                encoded = image_data
                mime_type = "image/png"

            # Decode and analyze the image
            image_bytes = base64.b64decode(encoded)
            original_size_mb = len(image_bytes) / (1024 * 1024)
            base64_size_mb = len(encoded) / (1024 * 1024)

            self.log.debug(
                f"Original image: {original_size_mb:.2f} MB (decoded), {base64_size_mb:.2f} MB (base64), type: {mime_type}"
            )

            # Determine optimization strategy
            reasons: List[str] = []
            if original_size_mb > max_size_mb:
                reasons.append(f"size > {max_size_mb} MB")
            if base64_size_mb > max_size_mb * 1.4:
                reasons.append("base64 overhead")
            if mime_type == "image/png" and original_size_mb > png_threshold:
                reasons.append(f"PNG > {png_threshold}MB")

            # Always check dimensions
            with Image.open(io.BytesIO(image_bytes)) as img:
                width, height = img.size
                resized_flag = False
                if width > max_dimension or height > max_dimension:
                    reasons.append(f"dimensions > {max_dimension}px")

                # Early exit: no optimization triggers -> keep original, record stats
                if not reasons:
                    if stats_list is not None:
                        stats_list.append(
                            {
                                "original_size_mb": round(original_size_mb, 4),
                                "final_size_mb": round(original_size_mb, 4),
                                "quality": None,
                                "format": mime_type.split("/")[-1].upper(),
                                "resized": False,
                                "reasons": ["no_optimization_needed"],
                                "final_hash": hashlib.sha256(
                                    encoded.encode()
                                ).hexdigest(),
                            }
                        )
                    self.log.debug(
                        "Skipping optimization: image already within thresholds"
                    )
                    return image_data

                self.log.debug(f"Optimization triggers: {', '.join(reasons)}")

                # Convert to RGB for JPEG compression
                if img.mode in ("RGBA", "LA", "P"):
                    background = Image.new("RGB", img.size, (255, 255, 255))
                    if img.mode == "P":
                        img = img.convert("RGBA")
                    background.paste(
                        img,
                        mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None,
                    )
                    img = background
                elif img.mode != "RGB":
                    img = img.convert("RGB")

                # Resize if needed
                if width > max_dimension or height > max_dimension:
                    ratio = min(max_dimension / width, max_dimension / height)
                    new_size = (int(width * ratio), int(height * ratio))
                    self.log.debug(
                        f"Resizing from {width}x{height} to {new_size[0]}x{new_size[1]}"
                    )
                    img = img.resize(new_size, Image.Resampling.LANCZOS)
                    resized_flag = True

                # Determine quality levels based on original size and user configuration
                if original_size_mb > 5.0:
                    quality_levels = [
                        base_quality,
                        base_quality - 10,
                        base_quality - 20,
                        base_quality - 30,
                        base_quality - 40,
                        max(base_quality - 50, 25),
                    ]
                elif original_size_mb > 2.0:
                    quality_levels = [
                        base_quality,
                        base_quality - 5,
                        base_quality - 15,
                        base_quality - 25,
                        max(base_quality - 35, 35),
                    ]
                else:
                    quality_levels = [
                        min(base_quality + 5, 95),
                        base_quality,
                        base_quality - 10,
                        max(base_quality - 20, 50),
                    ]

                # Ensure quality levels are within valid range (1-100)
                quality_levels = [max(1, min(100, q)) for q in quality_levels]

                # Try compression levels
                for quality in quality_levels:
                    output_buffer = io.BytesIO()
                    format_type = (
                        "JPEG"
                        if original_size_mb > png_threshold or "jpeg" in mime_type
                        else "PNG"
                    )
                    output_mime = f"image/{format_type.lower()}"

                    img.save(
                        output_buffer,
                        format=format_type,
                        quality=quality,
                        optimize=True,
                    )
                    output_bytes = output_buffer.getvalue()
                    output_size_mb = len(output_bytes) / (1024 * 1024)

                    if output_size_mb <= max_size_mb:
                        optimized_b64 = base64.b64encode(output_bytes).decode("utf-8")
                        self.log.debug(
                            f"Optimized: {original_size_mb:.2f} MB → {output_size_mb:.2f} MB (Q{quality})"
                        )
                        if stats_list is not None:
                            stats_list.append(
                                {
                                    "original_size_mb": round(original_size_mb, 4),
                                    "final_size_mb": round(output_size_mb, 4),
                                    "quality": quality,
                                    "format": format_type,
                                    "resized": resized_flag,
                                    "reasons": reasons,
                                    "final_hash": hashlib.sha256(
                                        optimized_b64.encode()
                                    ).hexdigest(),
                                }
                            )
                        return f"data:{output_mime};base64,{optimized_b64}"

                # Fallback: minimum quality
                output_buffer = io.BytesIO()
                img.save(output_buffer, format="JPEG", quality=15, optimize=True)
                output_bytes = output_buffer.getvalue()
                output_size_mb = len(output_bytes) / (1024 * 1024)
                optimized_b64 = base64.b64encode(output_bytes).decode("utf-8")

                self.log.warning(
                    f"Aggressive optimization: {output_size_mb:.2f} MB (Q15)"
                )
                if stats_list is not None:
                    stats_list.append(
                        {
                            "original_size_mb": round(original_size_mb, 4),
                            "final_size_mb": round(output_size_mb, 4),
                            "quality": 15,
                            "format": "JPEG",
                            "resized": resized_flag,
                            "reasons": reasons + ["fallback_min_quality"],
                            "final_hash": hashlib.sha256(
                                optimized_b64.encode()
                            ).hexdigest(),
                        }
                    )
                return f"data:image/jpeg;base64,{optimized_b64}"

        except Exception as e:
            self.log.error(f"Image optimization failed: {e}")
            # Return original or safe fallback
            if image_data.startswith("data:"):
                if stats_list is not None:
                    stats_list.append(
                        {
                            "original_size_mb": None,
                            "final_size_mb": None,
                            "quality": None,
                            "format": None,
                            "resized": False,
                            "reasons": ["optimization_failed"],
                            "final_hash": (
                                hashlib.sha256(encoded.encode()).hexdigest()
                                if "encoded" in locals()
                                else None
                            ),
                        }
                    )
                return image_data
            return f"data:image/jpeg;base64,{encoded if 'encoded' in locals() else image_data}"

    async def _fetch_file_as_base64(self, file_url: str) -> Optional[str]:
        """
        Fetch a file from Open WebUI's file system and convert to base64.

        Args:
            file_url: File URL from Open WebUI

        Returns:
            Base64 encoded file data or None if file not found
        """
        try:
            if "/api/v1/files/" in file_url:
                fid = file_url.split("/api/v1/files/")[-1].split("/")[0].split("?")[0]
            else:
                fid = file_url.split("/files/")[-1].split("/")[0].split("?")[0]

            from open_webui.models.files import Files

            file_obj = await Files.get_file_by_id(fid)
            if file_obj and file_obj.path:
                async with aiofiles.open(file_obj.path, "rb") as fp:
                    raw = await fp.read()
                enc = base64.b64encode(raw).decode()
                mime = file_obj.meta.get("content_type", "image/png")
                return f"data:{mime};base64,{enc}"
        except Exception as e:
            self.log.warning(f"Could not fetch file {file_url}: {e}")
        return None

    async def _upload_binary_file(
        self,
        __request__: Request,
        user: UserModel,
        file_bytes: bytes,
        filename: str,
        mime_type: str,
    ) -> str:
        bio = io.BytesIO(file_bytes)
        bio.seek(0)

        up_obj = await upload_file_handler(
            __request__,
            file=UploadFile(
                file=bio,
                filename=filename,
                headers=Headers({"content-type": mime_type}),
            ),
            metadata={"mime_type": mime_type, "source": "ppt_merge"},
            process=False,
            user=user,
        )

        return __request__.app.url_path_for("get_file_content_by_id", id=up_obj.id)

    async def _upload_image_with_status(
        self,
        image_data: Any,
        mime_type: str,
        __request__: Request,
        __user__: dict,
        __event_emitter__: Callable,
    ) -> str:
        """
        Unified image upload method with status updates and fallback handling.

        Returns:
            URL to uploaded image or data URL fallback
        """
        try:

            self.user = user = await Users.get_user_by_id(__user__["id"])

            # Convert image data to base64 string if needed
            if isinstance(image_data, bytes):
                image_data_b64 = base64.b64encode(image_data).decode("utf-8")
            else:
                image_data_b64 = str(image_data)

            image_url = await self._upload_image(
                __request__=__request__,
                user=user,
                image_data=image_data_b64,
                mime_type=mime_type,
            )

            # await __event_emitter__(
            #     {
            #         "type": "status",
            #         "data": {
            #             "action": "image_upload",
            #             "description": "图片存储成功!",
            #             "done": True,
            #         },
            #     }
            # )

            return image_url

        except Exception as e:
            self.log.warning(f"File upload failed, falling back to data URL: {e}")

            if isinstance(image_data, bytes):
                image_data_b64 = base64.b64encode(image_data).decode("utf-8")
            else:
                image_data_b64 = str(image_data)

            return f"data:{mime_type};base64,{image_data_b64}"

    async def _upload_image(
        self, __request__: Request, user: UserModel, image_data: str, mime_type: str
    ) -> str:
        """
        Upload generated image to Open WebUI's file system.
        Expects base64 encoded string input.

        Args:
            __request__: FastAPI request object
            user: User model object
            image_data: Base64 encoded image data string
            mime_type: MIME type of the image

        Returns:
            URL to the uploaded image or data URL fallback
        """
        try:
            self.log.debug(
                f"Processing image data, type: {type(image_data)}, length: {len(image_data)}"
            )

            # Decode base64 string to bytes
            try:
                decoded_data = base64.b64decode(image_data)
                self.log.debug(
                    f"Successfully decoded image data: {len(decoded_data)} bytes"
                )
            except Exception as decode_error:
                self.log.error(f"Failed to decode base64 data: {decode_error}")
                # Try to add padding if missing
                try:
                    missing_padding = len(image_data) % 4
                    if missing_padding:
                        image_data += "=" * (4 - missing_padding)
                    decoded_data = base64.b64decode(image_data)
                    self.log.debug(
                        f"Successfully decoded with padding: {len(decoded_data)} bytes"
                    )
                except Exception as second_decode_error:
                    self.log.error(f"Still failed to decode: {second_decode_error}")
                    return f"data:{mime_type};base64,{image_data}"

            bio = io.BytesIO(decoded_data)
            bio.seek(0)

            # Determine file extension
            extension = "png"
            if "jpeg" in mime_type or "jpg" in mime_type:
                extension = "jpg"
            elif "webp" in mime_type:
                extension = "webp"
            elif "gif" in mime_type:
                extension = "gif"

            # Create filename
            filename = f"gemini-generated-{uuid.uuid4().hex}.{extension}"

            # Upload with simple approach like reference
            up_obj = await upload_file_handler(
                __request__,
                file=UploadFile(
                    file=bio,
                    filename=filename,
                    headers=Headers({"content-type": mime_type}),
                ),
                metadata={"mime_type": mime_type, "source": "gemini_image_generation"},
                process=False,
                user=user,
            )

            self.log.debug(
                f"Upload completed. File ID: {up_obj.id}, Decoded size: {len(decoded_data)} bytes"
            )

            # Generate URL using reference method
            return __request__.app.url_path_for("get_file_content_by_id", id=up_obj.id)

        except Exception as e:
            self.log.exception(f"Image upload failed, using data URL fallback: {e}")
            # Fallback to data URL if upload fails
            return f"data:{mime_type};base64,{image_data}"

    def _configure_generation(
        self,
        body: Dict[str, Any],
        system_instruction: Optional[str],
        __metadata__: Dict[str, Any],
        __tools__: dict[str, Any] | None = None,
        enable_image_generation: bool = False,
        model_id: str = "",
    ) -> types.GenerateContentConfig:
        """
        Configure generation parameters and safety settings.

        Args:
            body: The request body containing generation parameters
            system_instruction: Optional system instruction string
            enable_image_generation: Whether to enable image generation
            model_id: The model ID being used (for feature support checks)

        Returns:
            types.GenerateContentConfig
        """
        gen_config_params = {
            "temperature": body.get("temperature"),
            "top_p": body.get("top_p"),
            "top_k": body.get("top_k"),
            "max_output_tokens": body.get("max_tokens"),
            "stop_sequences": body.get("stop") or None,
            "system_instruction": system_instruction,
        }

        # Enable image generation if requested
        if enable_image_generation:
            gen_config_params["response_modalities"] = ["TEXT", "IMAGE"]
            # 配置宽高比和分辨率
            gen_config_params["image_config"] = types.ImageConfig(
                aspect_ratio=self.valves.ASPECT_RATIO, image_size=self.valves.RESOLUTION
            )

        # Configure Gemini thinking/reasoning for models that support it
        # This is independent of include_thoughts - thinking config controls HOW the model reasons,
        # while include_thoughts controls whether the reasoning is shown in the output
        if self._check_thinking_support(model_id):
            try:
                thinking_config_params: Dict[str, Any] = {}

                # Determine include_thoughts setting
                include_thoughts = body.get("include_thoughts", True)
                if not self.valves.INCLUDE_THOUGHTS:
                    include_thoughts = False
                    self.log.debug(
                        "Thoughts output disabled via GOOGLE_INCLUDE_THOUGHTS"
                    )
                thinking_config_params["include_thoughts"] = include_thoughts

                # Check if model supports thinking_level (Gemini 3 models)
                if self._check_thinking_level_support(model_id):
                    # For Gemini 3 models, use thinking_level (not thinking_budget)
                    validated_level = self._validate_thinking_level(
                        self.valves.THINKING_LEVEL
                    )
                    if validated_level:
                        thinking_config_params["thinking_level"] = validated_level
                        self.log.debug(
                            f"Using thinking_level='{validated_level}' for model {model_id}"
                        )
                    else:
                        self.log.debug(
                            f"Using default thinking level for model {model_id}"
                        )
                else:
                    # For non-Gemini 3 models (e.g., Gemini 2.5), use thinking_budget
                    validated_budget = self._validate_thinking_budget(
                        self.valves.THINKING_BUDGET
                    )
                    if validated_budget == 0:
                        # Disable thinking if budget is 0
                        thinking_config_params["thinking_budget"] = 0
                        self.log.debug(
                            f"Thinking disabled via thinking_budget=0 for model {model_id}"
                        )
                    elif validated_budget > 0:
                        thinking_config_params["thinking_budget"] = validated_budget
                        self.log.debug(
                            f"Using thinking_budget={validated_budget} for model {model_id}"
                        )
                    else:
                        # -1 means dynamic thinking
                        thinking_config_params["thinking_budget"] = -1
                        self.log.debug(
                            f"Using dynamic thinking (model decides) for model {model_id}"
                        )

                gen_config_params["thinking_config"] = types.ThinkingConfig(
                    **thinking_config_params
                )
            except (AttributeError, TypeError) as e:
                # Fall back if SDK/model does not support ThinkingConfig
                self.log.debug(f"ThinkingConfig not supported: {e}")
            except Exception as e:
                # Log unexpected errors but continue without thinking config
                self.log.warning(f"Unexpected error configuring ThinkingConfig: {e}")

        # Configure safety settings
        if self.valves.USE_PERMISSIVE_SAFETY:
            safety_settings = [
                types.SafetySetting(
                    category="HARM_CATEGORY_HARASSMENT", threshold="BLOCK_NONE"
                ),
                types.SafetySetting(
                    category="HARM_CATEGORY_HATE_SPEECH", threshold="BLOCK_NONE"
                ),
                types.SafetySetting(
                    category="HARM_CATEGORY_SEXUALLY_EXPLICIT", threshold="BLOCK_NONE"
                ),
                types.SafetySetting(
                    category="HARM_CATEGORY_DANGEROUS_CONTENT", threshold="BLOCK_NONE"
                ),
            ]
            gen_config_params |= {"safety_settings": safety_settings}

        features = __metadata__.get("features", {})
        if (
            features.get("google_search_tool", False)
            or self.valves.ENABLE_GOOGLE_SEARCH
        ):
            self.log.debug("Enabling Google search grounding")
            gen_config_params.setdefault("tools", []).append(
                types.Tool(google_search=types.GoogleSearch())
            )

        params = __metadata__.get("params", {})
        if features.get("vertex_ai_search", False) or (
            self.valves.USE_VERTEX_AI
            and (self.valves.VERTEX_AI_RAG_STORE or os.getenv("VERTEX_AI_RAG_STORE"))
        ):
            vertex_rag_store = (
                params.get("vertex_rag_store")
                or self.valves.VERTEX_AI_RAG_STORE
                or os.getenv("VERTEX_AI_RAG_STORE")
            )
            if vertex_rag_store:
                self.log.debug(
                    f"Enabling Vertex AI Search grounding: {vertex_rag_store}"
                )
                gen_config_params.setdefault("tools", []).append(
                    types.Tool(
                        retrieval=types.Retrieval(
                            vertex_ai_search=types.VertexAISearch(
                                datastore=vertex_rag_store
                            )
                        )
                    )
                )
            else:
                self.log.warning(
                    "Vertex AI Search requested but vertex_rag_store not provided in params, valves, or env"
                )
        if __tools__ is not None and params.get("function_calling") == "native":
            for name, tool_def in __tools__.items():
                if not name.startswith("_"):
                    tool = tool_def["callable"]
                    self.log.debug(
                        f"Adding tool '{name}' with signature {tool.__signature__}"
                    )
                    gen_config_params.setdefault("tools", []).append(tool)

        # Filter out None values for generation config
        filtered_params = {k: v for k, v in gen_config_params.items() if v is not None}
        return types.GenerateContentConfig(**filtered_params)

    @staticmethod
    def _format_grounding_chunks_as_sources(
        grounding_chunks: list[types.GroundingChunk],
    ):
        formatted_sources = []
        for chunk in grounding_chunks:
            if hasattr(chunk, "retrieved_context") and chunk.retrieved_context:
                context = chunk.retrieved_context
                formatted_sources.append(
                    {
                        "source": {
                            "name": getattr(context, "title", None) or "Document",
                            "type": "vertex_ai_search",
                            "uri": getattr(context, "uri", None),
                        },
                        "document": [getattr(context, "chunk_text", None) or ""],
                        "metadata": [
                            {"source": getattr(context, "title", None) or "Document"}
                        ],
                    }
                )
            elif hasattr(chunk, "web") and chunk.web:
                context = chunk.web
                uri = context.uri
                title = context.title or "Source"

                formatted_sources.append(
                    {
                        "source": {
                            "name": title,
                            "type": "web_search_results",
                            "url": uri,
                        },
                        "document": ["Click the link to view the content."],
                        "metadata": [{"source": title}],
                    }
                )
        return formatted_sources

    async def _process_grounding_metadata(
        self,
        grounding_metadata_list: List[types.GroundingMetadata],
        text: str,
        __event_emitter__: Callable,
        *,
        emit_replace: bool = True,
    ):
        """Process and emit grounding metadata events."""
        grounding_chunks = []
        web_search_queries = []
        grounding_supports = []

        for metadata in grounding_metadata_list:
            if metadata.grounding_chunks:
                grounding_chunks.extend(metadata.grounding_chunks)
            if metadata.web_search_queries:
                web_search_queries.extend(metadata.web_search_queries)
            if metadata.grounding_supports:
                grounding_supports.extend(metadata.grounding_supports)

        # Add sources to the response
        if grounding_chunks:
            sources = self._format_grounding_chunks_as_sources(grounding_chunks)
            await __event_emitter__(
                {"type": "chat:completion", "data": {"sources": sources}}
            )

        # Add status specifying google queries used for grounding
        if web_search_queries:
            await __event_emitter__(
                {
                    "type": "status",
                    "data": {
                        "action": "web_search",
                        "description": "此回复基于 Google 搜索结果生成",
                        "urls": [
                            f"https://www.google.com/search?q={query}"
                            for query in web_search_queries
                        ],
                    },
                }
            )

        # Add citations in the text body
        replaced_text: Optional[str] = None
        if grounding_supports:
            # Citation indexes are in bytes
            ENCODING = "utf-8"
            text_bytes = text.encode(ENCODING)
            last_byte_index = 0
            cited_chunks = []

            for support in grounding_supports:
                cited_chunks.append(
                    text_bytes[last_byte_index : support.segment.end_index].decode(
                        ENCODING
                    )
                )

                # Generate and append citations (e.g., "[1][2]")
                footnotes = "".join(
                    [f"[{i + 1}]" for i in support.grounding_chunk_indices]
                )
                cited_chunks.append(f" {footnotes}")

                # Update index for the next segment
                last_byte_index = support.segment.end_index

            # Append any remaining text after the last citation
            if last_byte_index < len(text_bytes):
                cited_chunks.append(text_bytes[last_byte_index:].decode(ENCODING))

            replaced_text = "".join(cited_chunks)
            if emit_replace:
                await __event_emitter__(
                    {
                        "type": "replace",
                        "data": {"content": replaced_text},
                    }
                )

        # Return the transformed text when requested by caller
        if not emit_replace:
            return replaced_text if replaced_text is not None else text

    async def _handle_streaming_response(
        self,
        response_iterator: Any,
        __event_emitter__: Callable,
        __request__: Optional[Request] = None,
        __user__: Optional[dict] = None,
    ) -> AsyncIterator[str]:
        """
        Handle streaming response from Gemini API.

        Args:
            response_iterator: Iterator from generate_content
            __event_emitter__: Event emitter for status updates

        Returns:
            Generator yielding text chunks
        """

        async def emit_chat_event(event_type: str, data: Dict[str, Any]) -> None:
            if not __event_emitter__:
                return
            try:
                await __event_emitter__({"type": event_type, "data": data})
            except Exception as emit_error:  # pragma: no cover - defensive
                self.log.warning(f"Failed to emit {event_type} event: {emit_error}")

        await emit_chat_event("chat:start", {"role": "assistant"})

        grounding_metadata_list = []
        # Accumulate content separately for answer and thoughts
        answer_chunks: list[str] = []
        thought_chunks: list[str] = []
        thinking_started_at: Optional[float] = None

        try:
            async for chunk in response_iterator:
                # Check for safety feedback or empty chunks
                if not chunk.candidates:
                    # Check prompt feedback
                    if (
                        response_iterator.prompt_feedback
                        and response_iterator.prompt_feedback.block_reason
                    ):
                        block_reason = (
                            response_iterator.prompt_feedback.block_reason.name
                        )
                        message = f"[Blocked due to Prompt Safety: {block_reason}]"
                        await emit_chat_event(
                            "chat:finish",
                            {
                                "role": "assistant",
                                "content": message,
                                "done": True,
                                "error": True,
                            },
                        )
                        yield message
                    else:
                        message = "[Blocked by safety settings]"
                        await emit_chat_event(
                            "chat:finish",
                            {
                                "role": "assistant",
                                "content": message,
                                "done": True,
                                "error": True,
                            },
                        )
                        yield message
                    return  # Stop generation

                if chunk.candidates[0].grounding_metadata:
                    grounding_metadata_list.append(
                        chunk.candidates[0].grounding_metadata
                    )
                # Prefer fine-grained parts to split thoughts vs. normal text
                parts = []
                try:
                    parts = chunk.candidates[0].content.parts or []
                except Exception as parts_error:
                    # Fallback: use aggregated text if parts aren't accessible
                    self.log.warning(f"Failed to access content parts: {parts_error}")
                    if hasattr(chunk, "text") and chunk.text:
                        answer_chunks.append(chunk.text)
                        await __event_emitter__(
                            {
                                "type": "chat:message:delta",
                                "data": {
                                    "role": "assistant",
                                    "content": chunk.text,
                                },
                            }
                        )
                    continue

                for part in parts:
                    try:
                        # Thought parts (internal reasoning)
                        if getattr(part, "thought", False) and getattr(
                            part, "text", None
                        ):
                            if thinking_started_at is None:
                                thinking_started_at = time.time()
                            thought_chunks.append(part.text)
                            # Emit a live preview of what is currently being thought
                            preview = part.text.replace("\n", " ").strip()
                            MAX_PREVIEW = 120
                            if len(preview) > MAX_PREVIEW:
                                preview = preview[:MAX_PREVIEW].rstrip() + "…"
                            await __event_emitter__(
                                {
                                    "type": "status",
                                    "data": {
                                        "action": "thinking",
                                        "description": f"Thinking… {preview}",
                                        "done": False,
                                        "hidden": False,
                                    },
                                }
                            )

                        # Regular answer text
                        elif getattr(part, "text", None):
                            answer_chunks.append(part.text)
                            await __event_emitter__(
                                {
                                    "type": "chat:message:delta",
                                    "data": {
                                        "role": "assistant",
                                        "content": part.text,
                                    },
                                }
                            )
                    except Exception as part_error:
                        # Log part processing errors but continue with the stream
                        self.log.warning(f"Error processing content part: {part_error}")
                        continue

            # After processing all chunks, handle grounding data
            final_answer_text = "".join(answer_chunks)
            if grounding_metadata_list and __event_emitter__:
                # Don't emit replace here; we'll compose final content below
                cited = await self._process_grounding_metadata(
                    grounding_metadata_list,
                    final_answer_text,
                    __event_emitter__,
                    emit_replace=False,
                )
                final_answer_text = cited or final_answer_text

            final_content = final_answer_text
            details_block: Optional[str] = None

            if thought_chunks:
                duration_s = int(
                    max(0, time.time() - (thinking_started_at or time.time()))
                )
                # Format each line with > for blockquote while preserving formatting
                thought_content = "".join(thought_chunks).strip()
                quoted_lines = []
                for line in thought_content.split("\n"):
                    quoted_lines.append(f"> {line}")
                quoted_content = "\n".join(quoted_lines)

                details_block = f"""<details>
<summary>Thought ({duration_s}s)</summary>

{quoted_content}

</details>""".strip()
                final_content = f"{details_block}{final_answer_text}"

            if not final_content:
                final_content = ""

            # Ensure downstream consumers (UI, TTS) receive the complete response once streaming ends.
            await emit_chat_event(
                "replace", {"role": "assistant", "content": final_content}
            )
            await emit_chat_event(
                "chat:message",
                {"role": "assistant", "content": final_content, "done": True},
            )

            if thought_chunks:
                # Clear the thinking status without a summary in the status emitter
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {"action": "thinking", "done": True, "hidden": True},
                    }
                )

            await emit_chat_event(
                "chat:finish",
                {"role": "assistant", "content": final_content, "done": True},
            )

        except Exception as e:
            self.log.exception(f"Error during streaming: {e}")
            # Check if it's a chunk size error and provide specific guidance
            error_msg = str(e).lower()
            if "chunk too big" in error_msg or "chunk size" in error_msg:
                message = "Error: Image too large for processing. Please try with a smaller image (max 15 MB recommended) or reduce image quality."
            elif "quota" in error_msg or "rate limit" in error_msg:
                message = "Error: API quota exceeded. Please try again later."
            else:
                message = f"Error during streaming: {e}"
            await emit_chat_event(
                "chat:finish",
                {
                    "role": "assistant",
                    "content": message,
                    "done": True,
                    "error": True,
                },
            )
            yield message

    def _get_safety_block_message(self, response: Any) -> Optional[str]:
        """Check for safety blocks and return appropriate message."""
        # Check prompt feedback
        if response.prompt_feedback and response.prompt_feedback.block_reason:
            return f"[Blocked due to Prompt Safety: {response.prompt_feedback.block_reason.name}]"

        # Check candidates
        if not response.candidates:
            return "[Blocked by safety settings or no candidates generated]"

        # Check candidate finish reason
        candidate = response.candidates[0]
        if candidate.finish_reason == types.FinishReason.SAFETY:
            blocking_rating = next(
                (r for r in candidate.safety_ratings if r.blocked), None
            )
            reason = f" ({blocking_rating.category.name})" if blocking_rating else ""
            return f"[Blocked by safety settings{reason}]"
        elif candidate.finish_reason == types.FinishReason.PROHIBITED_CONTENT:
            return "[Content blocked due to prohibited content policy violation]"

        return None

    async def _retry_with_backoff(self, func, *args, **kwargs) -> Any:
        """
        Retry a function with exponential backoff.

        Args:
            func: Async function to retry
            *args, **kwargs: Arguments to pass to the function

        Returns:
            Result from the function

        Raises:
            The last exception encountered after all retries
        """
        max_retries = 3
        retry_count = 0
        last_exception = None

        while retry_count <= max_retries:
            try:
                return await func(*args, **kwargs)
            except ServerError as e:
                # These errors might be temporary, so retry
                retry_count += 1
                last_exception = e

                if retry_count <= max_retries:
                    # Calculate backoff time (exponential with jitter)
                    wait_time = min(2**retry_count + (0.1 * retry_count), 10)
                    self.log.warning(
                        f"Temporary error from Google API: {e}. Retrying in {wait_time:.1f}s ({retry_count}/{max_retries})"
                    )
                    await asyncio.sleep(wait_time)
                else:
                    raise
            except Exception:
                # Don't retry other exceptions
                raise

        # If we get here, we've exhausted retries
        assert last_exception is not None
        raise last_exception

    async def _images_to_ppt_bytes(self, images: List[bytes]) -> bytes:
        prs = Presentation()
        blank = prs.slide_layouts[6]

        for img in images:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                io.BytesIO(img),
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()

    async def _images_to_pdf_bytes(self, images: List[bytes]) -> bytes:
        pil_images = [Image.open(io.BytesIO(img)).convert("RGB") for img in images]

        output = io.BytesIO()
        pil_images[0].save(
            output,
            format="PDF",
            save_all=True,
            append_images=pil_images[1:],
        )
        output.seek(0)
        return output.read()

    async def _images_to_editable_ppt_bytes_keep_resolution(
        self, images: list[bytes], titles: list[str] = None
    ) -> bytes:
        """
        Convert images to editable PPT:
        - Each slide has a title textbox (editable)
        - Image inserted in original resolution (no scaling)
        """
        prs = Presentation()
        blank = prs.slide_layouts[6]

        if titles is None:
            titles = [f"Slide {i+1}" for i in range(len(images))]

        for img_bytes, title in zip(images, titles):
            slide = prs.slides.add_slide(blank)

            # 添加可编辑标题
            textbox = slide.shapes.add_textbox(
                Cm(1), Cm(0.5), width=prs.slide_width - Cm(2), height=Cm(2)
            )
            text_frame = textbox.text_frame
            text_frame.text = title
            text_frame.paragraphs[0].font.size = Pt(32)
            text_frame.paragraphs[0].font.bold = True

            # 插入图片，使用原始像素大小
            im = Image.open(io.BytesIO(img_bytes))
            width_px, height_px = im.size
            # PPTX 使用 EMU 单位 (1 px ≈ 9525 EMU)
            EMU_PER_PX = 9525
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                left=0,
                top=Cm(2.5),  # 从标题下方开始
                width=width_px * EMU_PER_PX,
                height=height_px * EMU_PER_PX,
            )

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()

    import io
    from PIL import Image
    from pptx import Presentation

    async def _images_to_editable_ppt_bytes_adaptive_no_distortion(
        self, images: list[bytes]
    ) -> bytes:
        """
        Convert images to editable PPT:
        - 移除标题 (No titles)
        - 自适应宽高且不变形 (Aspect Fit):
          图片会自动缩放以完整显示在幻灯片中，保持原始比例，并自动居中。
        """
        prs = Presentation()
        # 使用空白布局 (Layout 6 是空白页)
        blank_layout = prs.slide_layouts[6]

        # 获取幻灯片页面的物理尺寸 (EMU单位)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide_ratio = slide_width / slide_height

        for img_bytes in images:
            slide = prs.slides.add_slide(blank_layout)

            # 1. 读取图片信息
            # 注意：这里只读取尺寸，不改变图片本身，确保画质
            with Image.open(io.BytesIO(img_bytes)) as im:
                width_px, height_px = im.size
                img_ratio = width_px / height_px

            # 2. 计算目标尺寸和位置（核心算法：保持比例）
            if img_ratio > slide_ratio:
                # Case A: 图片比幻灯片更“扁/宽”
                # 策略：宽度撑满幻灯片，高度按比例缩小，垂直居中
                new_width = slide_width
                new_height = int(slide_width / img_ratio)

                left = 0
                top = (slide_height - new_height) // 2  # 垂直居中
            else:
                # Case B: 图片比幻灯片更“瘦/高” (或者比例相同)
                # 策略：高度撑满幻灯片，宽度按比例缩小，水平居中
                new_height = slide_height
                new_width = int(slide_height * img_ratio)

                top = 0
                left = (slide_width - new_width) // 2  # 水平居中

            # 3. 插入图片
            # 这里的 width 和 height 是严格根据原始比例计算的，所以绝对不会变形
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                left=left,
                top=top,
                width=new_width,
                height=new_height,
            )

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()

    async def _images_to_editable_ppt_bytes_fill_16_9(
        self, images: list[bytes]
    ) -> bytes:
        """
        Convert images to editable PPT:
        - 固定 16:9
        - 不变形
        - 无留白（Aspect Fill，必要时自动裁剪）
        """
        prs = Presentation()

        # 强制 16:9
        prs.slide_width = Inches(13.3333333333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide_ratio = slide_width / slide_height

        for img_bytes in images:
            slide = prs.slides.add_slide(blank_layout)

            with Image.open(io.BytesIO(img_bytes)) as im:
                img_w, img_h = im.size
                img_ratio = img_w / img_h

            # Aspect Fill：保证铺满
            if img_ratio > slide_ratio:
                # 图片更宽 → 高度撑满
                new_height = slide_height
                new_width = int(slide_height * img_ratio)

                top = 0
                left = (slide_width - new_width) // 2  # 负值 → 左右裁剪
            else:
                # 图片更高 → 宽度撑满
                new_width = slide_width
                new_height = int(slide_width / img_ratio)

                left = 0
                top = (slide_height - new_height) // 2  # 负值 → 上下裁剪

            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                left=left,
                top=top,
                width=new_width,
                height=new_height,
            )

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()

    async def _images_to_pdf_bytes_keep_resolution(self, images: list[bytes]) -> bytes:
        """
        Convert images to PDF with original resolution
        """
        pil_images = [Image.open(io.BytesIO(img)).convert("RGB") for img in images]

        # PDF 每页大小根据图片原始尺寸
        pdf_pages = []
        for im in pil_images:
            pdf_pages.append(im)

        output = io.BytesIO()
        pdf_pages[0].save(
            output, format="PDF", save_all=True, append_images=pdf_pages[1:]
        )
        output.seek(0)
        return output.read()

        import re  # 记得导入正则模块

    def _extract_dify_id_from_history(self, messages: list) -> str:
        """从历史消息中提取隐藏的 Dify Conversation ID"""
        # 倒序遍历，找最近的一条包含 ID 的助手消息
        # for msg in reversed(messages):
        #     if msg.get("role") == "assistant":
        #         content = msg.get("content", "")
        #         # 正则匹配 <!-- DIFY_ID:xxxx -->
        #         match = re.search(r"<!-- DIFY_ID:([a-zA-Z0-9-]+) -->", content)
        #         if match:
        #             return match.group(1)
        # return ""

        # for msg in reversed(messages):
        #     if msg.get("role") == "assistant":
        #         content = msg.get("content", "")

        #         # 正则匹配 [dify_id]: http://context/xxxx
        #         match = re.search(
        #             r"\[dify_id\]:\s*http://context/([a-zA-Z0-9-]+)", content
        #         )
        #         if match:
        #             return match.group(1)
        # return ""

        dify_id = ""
        processed_files = set()

        # 解释：
        # 1. \[dify_id\]: 匹配标签头
        # 2. ([a-zA-Z0-9-]+): 提取 Dify UUID (Group 1)
        # 3. \?files=([a-zA-Z0-9,]*): 提取文件列表 (Group 2)，用 * 号表示允许为空
        pattern = (
            r"\[dify_id\]:\s*http://context/([a-zA-Z0-9-]+)\?files=([a-zA-Z0-9,-]*)"
        )

        for msg in reversed(messages):
            if msg.get("role") == "assistant":
                content = msg.get("content", "")

                # 搜索匹配
                match = re.search(pattern, content)
                if match:
                    # 提取 Group 1: Dify ID
                    if not dify_id:
                        dify_id = match.group(1)

                    # 提取 Group 2: 文件列表 (即使为空也能匹配到空字符串)
                    file_ids_str = match.group(2)
                    if file_ids_str:
                        # 过滤掉可能的空字符串（比如 files=,, 的情况）
                        ids = [fid for fid in file_ids_str.split(",") if fid]
                        processed_files.update(ids)

                    # 找到一个就可以 break 了吗？
                    # 如果为了保险起见，想找全所有历史的文件ID，可以不 break，继续往前找
                    # 但通常只需要最新的那个状态就够了
                    break

        return dify_id, processed_files

    async def _upload_to_dify_direct(
        self,
        image_data: Union[List[dict], Any],  # 参数名保持不变，但逻辑上支持任意文件
        user_id: str,
        api_base: str,
        api_key: str,
        mime_type: str = None,
        max_retries: int = 3,
        initial_delay: float = 1.0,
    ) -> Union[List[dict], dict]:
        """
        [异步并发 + 自动重试版本]
        将图片或文件上传到 Dify。支持并发上传，并包含失败重试机制。
        """
        upload_url = f"{api_base}/v1/files/upload"
        headers = {"Authorization": f"Bearer {api_key}"}

        # === 内部函数：带重试逻辑的单文件上传 ===
        async def _upload_single_file_with_retry(session, current_data, current_mime):
            # 1. 数据预处理 (只做一次)
            file_content = None
            try:
                if isinstance(current_data, bytes):
                    file_content = current_data
                else:
                    img_str = str(current_data)
                    if "base64," in img_str:
                        img_str = img_str.split("base64,")[1]
                    try:
                        file_content = base64.b64decode(img_str)
                    except:
                        file_content = img_str.encode("utf-8")
            except Exception as e:
                print(f"[Dify] Data decode error, skipping: {e}")
                raise e

            # === [修改点 1] 优化文件后缀生成逻辑 ===
            # 原代码: ext = current_mime.split("/")[-1] if current_mime else "png"
            ext = "bin"  # 默认值
            if current_mime:
                guess_ext = mimetypes.guess_extension(current_mime)
                if guess_ext:
                    ext = guess_ext.lstrip(".")  # .pdf -> pdf
                else:
                    ext = current_mime.split("/")[-1]  # 降级处理

            filename = f"dify_upload_{uuid.uuid4().hex[:8]}.{ext}"

            # 2. 进入重试循环
            last_exception = None
            for attempt in range(max_retries + 1):
                try:
                    form_data = aiohttp.FormData()
                    form_data.add_field(
                        "file",
                        file_content,
                        filename=filename,
                        content_type=current_mime,
                    )
                    form_data.add_field("user", user_id)

                    async with session.post(
                        upload_url, headers=headers, data=form_data
                    ) as response:
                        if response.status in [200, 201]:
                            resp_text = await response.text()
                            result = json.loads(resp_text)

                            # === [修改点 2] 动态判断返回类型 ===
                            # 原代码: "type": "image",
                            is_image = current_mime and current_mime.startswith(
                                "image/"
                            )
                            res_type = "document"
                            if is_image:
                                res_type = "image"
                            else:
                                # 注意: 此列表需与 Dify DOCUMENT_EXTENSIONS (ETL_TYPE=dify) 保持一致
                                # pptx/ppt/xml/doc 在 ETL_TYPE=dify 时不被识别为 document，
                                # 发送 "document" 会触发 strict_type_validation 报错
                                if ext.lower() not in [
                                    "txt",
                                    "md",
                                    "markdown",
                                    "mdx",
                                    "pdf",
                                    "html",
                                    "htm",
                                    "xlsx",
                                    "xls",
                                    "vtt",
                                    "properties",
                                    "docx",
                                    "csv",
                                    "eml",
                                    "msg",
                                    "epub",
                                ]:
                                    res_type = "custom"
                            # res_type = "image" if is_image else "document"

                            return {
                                "type": res_type,
                                "transfer_method": "local_file",
                                "upload_file_id": result.get("id"),
                            }

                        elif 400 <= response.status < 500 and response.status != 429:
                            error_text = await response.text()
                            raise ValueError(
                                f"Dify Error {response.status}: {error_text}"
                            )
                        else:
                            error_text = await response.text()
                            msg = f"HTTP {response.status} - {error_text}"
                            raise aiohttp.ClientResponseError(
                                response.request_info,
                                response.history,
                                status=response.status,
                                message=msg,
                            )
                except (aiohttp.ClientError, asyncio.TimeoutError, ValueError) as e:
                    last_exception = e
                    if isinstance(e, ValueError):
                        raise e
                    if attempt < max_retries:
                        # 修正了原代码中的语法错误: initial_delay _(2_*attempt) -> * (2**attempt)
                        sleep_time = initial_delay * (2**attempt)
                        print(
                            f"[Dify Upload] Attempt {attempt+1}/{max_retries} failed: {e}. Retrying in {sleep_time}s..."
                        )
                        await asyncio.sleep(sleep_time)
                    else:
                        print(f"[Dify Upload] All {max_retries+1} attempts failed.")

            if last_exception:
                raise last_exception
            raise Exception("Unknown upload failure")

        # === 主逻辑 ===
        timeout = aiohttp.ClientTimeout(total=60)
        async with aiohttp.ClientSession(timeout=timeout) as session:
            # 场景 1: 列表格式
            if isinstance(image_data, list):
                tasks = []
                for item in image_data:
                    inline = item.get("inline_data", {})
                    img_content = inline.get("data")
                    img_mime = inline.get("mime_type")
                    if img_content:
                        tasks.append(
                            _upload_single_file_with_retry(
                                session, img_content, img_mime
                            )
                        )
                if not tasks:
                    return []
                print(f"[Dify] Processing {len(tasks)} files with retry mechanism...")
                return await asyncio.gather(*tasks)
            # 场景 2: 单个数据
            else:
                if not mime_type:
                    mime_type = "image/png"  # 默认值，如果传文件最好显式指定
                return await _upload_single_file_with_retry(
                    session, image_data, mime_type
                )

    @staticmethod
    def _get_public_base_url(request: Request) -> str:
        """从请求头中提取公网 base URL（如 https://chat.focusmedia.cn）"""
        origin = request.headers.get("origin")
        if origin:
            return origin.rstrip("/")
        scheme = request.headers.get("x-forwarded-proto", request.url.scheme)
        host = request.headers.get("x-forwarded-host", request.headers.get("host", ""))
        return f"{scheme}://{host}" if host else ""

    def _build_ppt_viewer_html(
        self,
        slides: list,
        dl_url_pdf: str = "",
        dl_url_ppt: str = "",
        base_url: str = "",
    ) -> str:
        """
        构建 PowerPoint 风格的 PPT 查看器 HTML embed。
        左侧缩略图导航，右侧大图预览，顶部下载按钮。
        """
        slides_json = json.dumps(
            [{"title": s.get("title", "Slide"), "url": s["url"]} for s in slides],
            ensure_ascii=False,
        )
        pdf_json = json.dumps(dl_url_pdf)
        ppt_json = json.dumps(dl_url_ppt)
        return (
            "<html><head><style>"
            "*{margin:0;padding:0;box-sizing:border-box}"
            'body{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,sans-serif;'
            "background:#1e1e2e;color:#e0e0e0;overflow-x:hidden}"
            ".container{display:flex;width:100%;min-height:420px}"
            ".sidebar{width:160px;min-width:160px;background:#252526;"
            "border-right:1px solid #3a3a4a;overflow-y:auto;padding:8px 6px;"
            "max-height:540px}"
            ".sidebar::-webkit-scrollbar{width:4px}"
            ".sidebar::-webkit-scrollbar-thumb{background:#555;border-radius:2px}"
            ".thumb-item{cursor:pointer;margin-bottom:8px;border-radius:4px;"
            "border:2px solid transparent;overflow:hidden;transition:border-color .2s;position:relative}"
            ".thumb-item:hover{border-color:#666}"
            ".thumb-item.active{border-color:#ff9800}"
            ".thumb-item img{width:100%;display:block;background:#333;aspect-ratio:16/9;object-fit:contain}"
            ".thumb-num{position:absolute;top:2px;left:2px;background:rgba(0,0,0,.7);"
            "color:#ccc;font-size:10px;padding:1px 5px;border-radius:3px}"
            ".main{flex:1;display:flex;flex-direction:column;min-width:0}"
            ".toolbar{display:flex;align-items:center;justify-content:space-between;"
            "padding:8px 16px;background:#2d2d3d;border-bottom:1px solid #3a3a4a}"
            ".slide-counter{font-size:14px;color:#bbb}"
            ".dl-buttons{display:flex;gap:8px}"
            ".dl-btn{display:inline-flex;align-items:center;gap:4px;"
            "padding:5px 14px;border-radius:5px;font-size:12px;font-weight:600;"
            "cursor:pointer;text-decoration:none;border:none;transition:opacity .2s}"
            ".dl-btn:hover{opacity:.85}"
            ".dl-btn.pdf{background:#4a90d9;color:#fff}"
            ".dl-btn.pptx{background:#ff9800;color:#fff}"
            ".viewer{flex:1;display:flex;align-items:center;justify-content:center;"
            "position:relative;padding:12px;background:#1e1e2e}"
            ".viewer img{max-width:100%;max-height:100%;object-fit:contain;"
            "border-radius:4px;box-shadow:0 4px 20px rgba(0,0,0,.5)}"
            ".nav-btn{position:absolute;top:50%;transform:translateY(-50%);"
            "background:rgba(0,0,0,.5);color:#fff;border:none;width:36px;height:60px;"
            "font-size:24px;cursor:pointer;border-radius:4px;transition:background .2s;z-index:2}"
            ".nav-btn:hover{background:rgba(0,0,0,.8)}"
            ".nav-btn.prev{left:4px}"
            ".nav-btn.next{right:4px}"
            ".nav-btn:disabled{opacity:.3;cursor:default}"
            ".toast{position:fixed;bottom:20px;left:50%;transform:translateX(-50%);"
            "background:rgba(0,0,0,.75);color:#fff;padding:8px 20px;border-radius:8px;"
            "font-size:13px;z-index:99;pointer-events:none;opacity:0;transition:opacity .3s}"
            ".toast.show{opacity:1}"
            "@media(max-width:640px){"
            ".container{min-height:240px}"
            ".sidebar{width:72px;min-width:72px;padding:4px 3px;max-height:400px}"
            ".thumb-item{margin-bottom:4px}"
            ".thumb-num{font-size:8px;padding:0 3px}"
            ".toolbar{flex-wrap:wrap;gap:4px;padding:6px 8px}"
            ".dl-btn{padding:4px 8px;font-size:11px}"
            ".viewer{padding:4px}"
            ".nav-btn{width:26px;height:40px;font-size:16px}}"
            f'</style><base href="{base_url}"></head><body>'
            '<div class="container">'
            '<div class="sidebar" id="sidebar"></div>'
            '<div class="main">'
            '<div class="toolbar">'
            '<span class="slide-counter" id="counter">1 / 1</span>'
            '<div class="dl-buttons" id="dlbtns"></div>'
            "</div>"
            '<div class="viewer">'
            '<button class="nav-btn prev" id="prevBtn" onclick="go(-1)">\u2039</button>'
            '<img id="mainImg" src="" alt="slide"/>'
            '<button class="nav-btn next" id="nextBtn" onclick="go(1)">\u203a</button>'
            "</div></div></div>"
            '<div class="toast" id="toast"></div>'
            "<script>"
            "(function(){"
            f"var slides={slides_json},"
            f"pdfUrl={pdf_json},"
            f"pptUrl={ppt_json},"
            "cur=0,"
            'sb=document.getElementById("sidebar"),'
            'img=document.getElementById("mainImg"),'
            'ctr=document.getElementById("counter"),'
            'pb=document.getElementById("prevBtn"),'
            'nb=document.getElementById("nextBtn"),'
            'dlb=document.getElementById("dlbtns");'
            "function show(i){"
            "if(i<0||i>=slides.length)return;"
            "cur=i;img.src=slides[i].url;"
            'ctr.textContent=(i+1)+" / "+slides.length;'
            "pb.disabled=i===0;nb.disabled=i===slides.length-1;"
            'var thumbs=sb.querySelectorAll(".thumb-item");'
            "for(var j=0;j<thumbs.length;j++){"
            'thumbs[j].className="thumb-item"+(j===i?" active":"");'
            "}"
            'thumbs[i].scrollIntoView({block:"nearest",behavior:"smooth"});'
            "rh();"
            "}"
            "function go(d){show(cur+d);}"
            "window.go=go;"
            "for(var i=0;i<slides.length;i++){"
            '  var d=document.createElement("div");d.className="thumb-item";'
            '  var im=document.createElement("img");im.src=slides[i].url;im.alt="P"+(i+1);'
            '  im.loading=i<6?"eager":"lazy";'
            '  var n=document.createElement("span");n.className="thumb-num";n.textContent=i+1;'
            "  d.appendChild(im);d.appendChild(n);"
            '  d.setAttribute("data-i",i);'
            '  d.addEventListener("click",function(){show(parseInt(this.getAttribute("data-i")));});'
            "  sb.appendChild(d);"
            "}"
            "var _wx=/MicroMessenger|miniProgram/i.test(navigator.userAgent);"
            "function _toast(m){var t=document.getElementById('toast');if(!t)return;t.textContent=m;"
            "t.classList.add('show');setTimeout(function(){t.classList.remove('show');},2000);}"
            "function _copy(u){if(navigator.clipboard){navigator.clipboard.writeText(u)"
            ".then(function(){_toast('链接已复制，请在浏览器中打开下载');}).catch(function(){prompt('长按复制链接:',u);});"
            "}else{prompt('长按复制链接:',u);}}"
            "function dlClick(e){e.preventDefault();var u=this.href;"
            "if(_wx){_copy(u);return;}"
            "try{window.parent.open(u,'_blank');}catch(x){window.open(u,'_blank');}}"
            'if(pdfUrl){var a=document.createElement("a");a.className="dl-btn pdf";'
            'a.textContent="\u2b07 \u4e0b\u8f7dPDF";a.href=pdfUrl;a.onclick=dlClick;dlb.appendChild(a);}'
            'if(pptUrl){var a2=document.createElement("a");a2.className="dl-btn pptx";'
            'a2.textContent="\u2b07 \u4e0b\u8f7dPPTX";a2.href=pptUrl;a2.onclick=dlClick;dlb.appendChild(a2);}'
            "function rh(){"
            "var h=Math.max(document.documentElement.scrollHeight||0,document.body.scrollHeight||0);"
            'if(h>0)parent.postMessage({type:"iframe:height",height:h},"*");'
            "}"
            'document.addEventListener("keydown",function(e){'
            'if(e.key==="ArrowLeft"||e.key==="ArrowUp"){e.preventDefault();go(-1);}'
            'if(e.key==="ArrowRight"||e.key==="ArrowDown"){e.preventDefault();go(1);}'
            "});"
            "show(0);"
            'window.addEventListener("resize",rh);'
            "setTimeout(rh,100);setTimeout(rh,500);setTimeout(rh,2000);"
            "try{var _f=window.frameElement;if(_f){"
            "var _ec=_f.parentElement.parentElement,"
            "_cc=_ec.parentElement.querySelector('#response-content-container');"
            "if(_cc){var _mv=function(){try{if(_cc.nextElementSibling!==_ec)_cc.after(_ec);}catch(e){}};"
            "_mv();setInterval(_mv,2000);}"
            "}}catch(e){}"
            "})();"
            "</script>"
            "</body></html>"
        )

    def _build_ppt_progress_html(self, start_ts_ms: int = 0) -> str:
        """构建自更新 PPT 生成进度 HTML（embed iframe 内渲染）。"""
        if start_ts_ms <= 0:
            start_ts_ms = int(time.time() * 1000)
        return (
            "<html><head><style>"
            "*{margin:0;padding:0;box-sizing:border-box}"
            'body{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,sans-serif;overflow:hidden}'
            "@keyframes g{0%{background-position:0% 50%}50%{background-position:100% 50%}100%{background-position:0% 50%}}"
            "@keyframes pulse{0%,100%{opacity:1}50%{opacity:.6}}"
            ".wrap{position:relative;width:100%;border-radius:12px;overflow:hidden;"
            "background:linear-gradient(-45deg,#0d1b2a,#1b2838,#162447,#1a1a2e,#1f4068,#162447,#2d3a4a);"
            "background-size:400% 400%;animation:g 8s ease infinite;visibility:hidden;padding:20px}"
            ".overlay{position:absolute;inset:0;"
            "background:radial-gradient(ellipse at 30% 50%,rgba(255,152,0,.15),transparent 60%),"
            "radial-gradient(ellipse at 70% 60%,rgba(66,165,245,.15),transparent 50%)}"
            ".content{position:relative;display:flex;flex-direction:column;align-items:center;gap:10px}"
            ".badge{display:flex;align-items:center;gap:6px;"
            "background:rgba(0,0,0,.4);backdrop-filter:blur(8px);border-radius:8px;padding:5px 12px;"
            "font-size:13px;font-weight:600;color:#fff;align-self:flex-start}"
            ".icon-wrap{width:52px;height:52px;border-radius:50%;"
            "background:rgba(255,255,255,.08);backdrop-filter:blur(4px);"
            "display:flex;align-items:center;justify-content:center;"
            "animation:pulse 2s ease-in-out infinite}"
            ".icon-wrap svg{width:28px;height:28px}"
            ".title{color:#fff;font-size:15px;font-weight:600;text-shadow:0 1px 4px rgba(0,0,0,.4)}"
            ".tips-box{background:rgba(0,0,0,.3);backdrop-filter:blur(4px);border-radius:8px;"
            "padding:12px 14px;width:100%;max-width:520px;"
            "color:rgba(255,255,255,.75);font-size:12px;line-height:1.8;text-align:left}"
            ".tips-title{color:#ffb74d;font-weight:600;font-size:13px;margin-bottom:4px}"
            ".tips-box a{color:#64b5f6;text-decoration:none}"
            ".tips-box a:hover{text-decoration:underline}"
            ".bar-wrap{width:65%;max-width:340px;height:5px;"
            "background:rgba(255,255,255,.12);border-radius:3px;overflow:hidden}"
            ".bar{height:100%;border-radius:3px;transition:width .6s ease;"
            "background:linear-gradient(90deg,#ff9800,#42a5f5)}"
            ".time{color:rgba(255,255,255,.4);font-size:11px}"
            "</style></head><body>"
            '<div class="wrap"><div class="overlay"></div>'
            '<div class="content">'
            '<div class="badge"><span id="pct">0%</span> <span id="stxt">'
            "\u51c6\u5907\u4e2d</span></div>"
            '<div class="icon-wrap">'
            '<svg viewBox="0 0 24 24" fill="none" xmlns="http://www.w3.org/2000/svg">'
            '<rect x="3" y="3" width="18" height="18" rx="3" stroke="white" '
            'stroke-width="1.5" fill="rgba(255,152,0,.3)"/>'
            '<text x="12" y="16" text-anchor="middle" font-size="10" '
            'font-weight="bold" fill="white">PPT</text>'
            "</svg></div>"
            '<div class="title" id="ttl">&#9203; PPT \u6b63\u5728\u751f\u6210\u4e2d</div>'
            '<div class="tips-box">'
            '<div class="tips-title">\u4f7f\u7528\u6e29\u99a8\u63d0\u793a</div>'
            "1. \u6b63\u5728\u751f\u4ea7\u9884\u8ba13-10\u5206\u949f\uff0c"
            "\u60a8\u53ef\u5728\u6d88\u606f\u9876\u90e8\u67e5\u770b\u751f\u6210\u8fdb\u5ea6\uff0c"
            "\u8bf7\u52ff\u91cd\u590d\u70b9\u51fb\u751f\u6210\uff01<br>"
            "2. PPT\u4e3a\u5355\u4efb\u52a1\u6a21\u5f0f\uff0c\u4e0d\u652f\u6301\u591a\u8f6e\u5bf9\u8bdd\uff0c"
            "\u5982\u9700\u4fee\u6539\u8bf7\u91cd\u65b0\u63d0\u4ea4"
            "&quot;\u5b8c\u6574\u5927\u7eb2\u6216\u9644\u4ef6&quot;\uff1b<br>"
            "3. \u652f\u6301\u5728\u7ebf\u67e5\u770b\u5e76\u652f\u6301\u5bfc\u51faPDF\u548cPPT\u683c\u5f0f\uff0c"
            "\u683c\u5f0f\u4e3a\u3010\u56fe\u7247\u578b\u3011\uff0c\u5982\u9700\u4fee\u6539\u5efa\u8bae\u4e0b\u8f7dPDF\u4f7f\u7528"
            "&quot;WPS\u8f6c\u6362&quot;\uff0c\u7f51\u5740\uff1a"
            '<a href="https://aippt.wps.cn/aippt/convert-ppt/home?request_source=banner" '
            'target="_blank">aippt.wps.cn</a>\uff1b'
            "</div>"
            '<div class="bar-wrap"><div class="bar" id="bar" style="width:0%"></div></div>'
            '<div class="time" id="tm">\u5df2\u7b49\u5f85 0\u79d2</div>'
            "</div></div>"
            "<script>"
            "(function(){"
            f"var start={start_ts_ms},"
            'pct=document.getElementById("pct"),'
            'stxt=document.getElementById("stxt"),'
            'ttl=document.getElementById("ttl"),'
            'bar=document.getElementById("bar"),'
            'tm=document.getElementById("tm");'
            "function calc(s){"
            'if(s<10)return{p:Math.floor(s/2),t:"\u51c6\u5907\u4e2d"};'
            'if(s<60)return{p:5+Math.floor(10*(s-10)/50),t:"\u89e3\u6790\u5927\u7eb2\u4e2d"};'
            'if(s<180)return{p:15+Math.floor(35*(s-60)/120),t:"\u751f\u6210\u5e7b\u706f\u7247\u4e2d"};'
            'if(s<420)return{p:50+Math.floor(35*(s-180)/240),t:"\u6e32\u67d3\u7f8e\u5316\u4e2d"};'
            "var prog=85+Math.floor(10*(1-1/(1+(s-420)/180)));"
            'return{p:Math.min(95,prog),t:"\u5373\u5c06\u5b8c\u6210"};'
            "}"
            "var first=true;"
            "function upd(){"
            "var s=Math.floor((Date.now()-start)/1000),r=calc(s),"
            "m=Math.floor(s/60),sec=s%60,"
            'tStr=m>0?m+"\u5206"+sec+"\u79d2":sec+"\u79d2";'
            "stxt.textContent=r.t;"
            'ttl.innerHTML=(s<10?"&#9203;":"&#9889;")+" PPT "+r.t;'
            'if(first){bar.style.transition="none";}'
            'pct.textContent=r.p+"%";'
            'bar.style.width=r.p+"%";'
            'tm.textContent="\u5df2\u7b49\u5f85 "+tStr;'
            'if(first){bar.offsetWidth;bar.style.transition="width .6s ease";first=false;}'
            'document.querySelector(".wrap").style.visibility="visible";'
            "var h=Math.max(document.documentElement.scrollHeight||0,document.body.scrollHeight||0);"
            'if(h>0)parent.postMessage({type:"iframe:height",height:h},"*");'
            "}"
            "setInterval(upd,1000);"
            "upd();"
            "try{var _f=window.frameElement;if(_f){"
            "var _ec=_f.parentElement.parentElement,"
            "_cc=_ec.parentElement.querySelector('#response-content-container');"
            "if(_cc){var _mv=function(){try{if(_cc.nextElementSibling!==_ec)_cc.after(_ec);}catch(e){}};"
            "_mv();setInterval(_mv,2000);}"
            "}}catch(e){}"
            "})();"
            "</script>"
            "</body></html>"
        )

    async def _emit_ppt_progress_once(
        self,
        emitter,
        start_ts_ms: int = 0,
        chat_id: str = "",
        message_id: str = "",
    ):
        """发送一次 PPT 自更新进度 embed（内部 JS 自动计时更新）。
        通过 event_emitter 推送给前端实时显示，同时延迟写 DB 兜底。
        原因：pipe() 返回后前端 saveChatHandler 会用不含 embeds 的 history
        整体覆盖 DB，导致 socket handler 写入的进度条被清掉。
        延迟 2 秒后再写 DB，确保在 saveChatHandler 之后执行，
        用户切换页面再切回来时仍能从 DB 恢复进度条。
        """
        if not emitter:
            return
        html = self._build_ppt_progress_html(start_ts_ms)
        await emitter({"type": "embeds", "data": {"embeds": [html]}})

        if chat_id and message_id:

            async def _persist_progress():
                try:
                    await asyncio.sleep(2)
                    from open_webui.models.chats import Chats

                    await Chats.upsert_message_to_chat_by_id_and_message_id(
                        chat_id,
                        message_id,
                        {"embeds": [html]},
                    )
                except Exception:
                    pass

            asyncio.create_task(_persist_progress())

    async def _clear_ppt_progress_embed(
        self, event_emitter, chat_id: str = "", message_id: str = ""
    ):
        """清除 DB 中的进度条 embed。
        仅做直接 DB 写入清空 embeds，不发 embeds 事件。
        原因：socket handler 的 embeds 事件处理会 extend(追加) DB 中旧的 embeds，
        发送空 embeds 事件反而会把旧进度条从 DB 读出并写回，导致清除失效。
        后续的 viewer embed 事件会同时更新客户端显示和 DB 状态。
        （参考 seedance.py 的处理方式）
        """
        if chat_id and message_id:
            try:
                from open_webui.models.chats import Chats

                await Chats.upsert_message_to_chat_by_id_and_message_id(
                    chat_id,
                    message_id,
                    {"embeds": []},
                )
            except Exception as e:
                self.log.warning(f"Clear progress embed from DB failed: {e}")

    async def _run_background_task(
        self,
        dify_url: str,
        workflow_id: str,
        processed_ids_list: list,
        query: str,
        files: Union[List[dict], dict],
        request: Request,
        event_emitter: Callable,
        __user__=None,
        chat_id: str = "",
        message_id: str = "",
        session_id: str = "",
    ):
        """这是在后台运行的实际逻辑，不会阻塞 HTTP 请求"""
        try:
            user = await Users.get_user_by_id(__user__["id"])
            start_ts = int(time.time() * 1000)
            await self._emit_ppt_progress_once(
                event_emitter,
                start_ts,
                chat_id=chat_id,
                message_id=message_id,
            )
            headers = {
                "Authorization": f"Bearer {self.valves.API_KEY}",
                "Content-Type": "application/json",
            }
            payload = {
                "inputs": {"query": query},
                "response_mode": "streaming",
                "user": user.email,
            }
            if files:
                payload["files"] = files
            print(f"payload:{payload}")

            final_outputs = {}
            final_error = None
            elapsed_time = 0
            last_node_outputs = (
                {}
            )  # [FIX] 记录最后一个 node_finished 的 outputs，作为 workflow_finished 丢失时的保底
            received_events = []  # [FIX] 记录收到的所有事件类型，用于诊断

            # 执行指定版本 Workflow
            if workflow_id:
                dify_url = dify_url.replace("run", f"{workflow_id}/run")

            # 1. 连接 Dify (Streaming)，失败最多重试3次
            # [FIX] 显式拆分 timeout：read_timeout 设大以防 PPT 生成期间长时间无 SSE 事件导致连接被判定超时
            stream_timeout = httpx.Timeout(connect=30, read=1200, write=30, pool=30)
            max_retries = 3
            for _attempt in range(1, max_retries + 1):
                final_outputs = {}
                final_error = None
                last_node_outputs = {}
                try:
                    async with httpx.AsyncClient(
                        timeout=stream_timeout, http2=False
                    ) as client:

                        async with client.stream(
                            "POST", dify_url, headers=headers, json=payload
                        ) as response:
                            print(
                                f"response.status_code==============================================={response.status_code}"
                            )
                            if response.status_code != 200:

                                raise Exception(f"API Error {response.status_code}")

                            raw_line_count = 0
                            async for line in response.aiter_lines():
                                if not line or not line.startswith("data:"):
                                    raw_line_count += 1
                                    if line:
                                        print(
                                            f"[SSE] 非data行({raw_line_count}): {line[:500]}"
                                        )
                                    continue
                                try:
                                    # data = json.loads(line[6:])
                                    json_str = line[
                                        5:
                                    ].lstrip()  # 兼容 "data: " 和 "data:" 两种 SSE 格式
                                    data = json.loads(json_str)
                                    event = data.get("event")
                                    content = data.get("data", {})
                                    received_events.append(event)

                                    # # 获取中间节点 "输出 3"的结果，用于后续PPT大纲的显示
                                    if content.get("node_id", "") == "1766476581627":
                                        if content.get("status") == "succeeded":

                                            await event_emitter(
                                                {
                                                    "type": "message",
                                                    "data": {
                                                        "content": content.get(
                                                            "outputs", {}
                                                        ).get("text", "")
                                                    },
                                                }
                                            )

                                    if event == "node_started":
                                        await event_emitter(
                                            {
                                                "type": "status",
                                                "data": {
                                                    "description": f"正在执行: {content.get('title', '...')}",
                                                    "done": False,
                                                },
                                            }
                                        )
                                    # [FIX] 捕获 node_finished 的 outputs 作为保底，防止 workflow_finished 丢失时无数据  新增的
                                    elif event == "node_finished":
                                        node_outputs = content.get("outputs", {})
                                        if node_outputs and node_outputs.get(
                                            "ppt_urls"
                                        ):
                                            last_node_outputs = node_outputs
                                            print(
                                                f"[FIX] node_finished 捕获到 ppt_urls, node_id={content.get('node_id')}"
                                            )
                                    elif event == "workflow_finished":
                                        print(
                                            f"workflow_finished raw data===============: {json.dumps(data, ensure_ascii=False, default=str)[:2000]}"
                                        )  # [FIX] 打印原始数据便于排查
                                        if content.get("status") == "succeeded":
                                            final_outputs = content.get("outputs", {})
                                            elapsed_time = int(
                                                content.get("elapsed_time", 0)
                                            )
                                        else:
                                            # [FIX] 原代码 content.get("error") 可能为 None，导致 final_error 为 None 跳过错误检查
                                            final_error = (
                                                content.get("error")
                                                or f"workflow status: {content.get('status')}"
                                            )

                                    elif event == "error":
                                        final_error = data.get("message")
                                except Exception as e:
                                    self.log.error(f"SSE Task Error: {str(e)}")
                                    print(
                                        f"SSE Task Error=================================================: {str(e)}"
                                    )
                                    print(
                                        f"SSE raw line===: {line[:500]}"
                                    )  # [FIX] 打印出错的原始 SSE 行便于排查

                    if final_error:
                        raise Exception(final_error)

                    # [FIX] 当 workflow_finished 未收到或 outputs 为空时，用 node_finished 中捕获的 outputs 保底
                    if not final_outputs and last_node_outputs:
                        print(
                            f"[FIX] workflow_finished 未收到或 outputs 为空，使用 node_finished 保底数据"
                        )
                        final_outputs = last_node_outputs
                    if not final_outputs:
                        raise Exception(
                            f"未收到 PPT 数据, 收到的事件列表: {received_events}"
                        )
                    break  # 成功，跳出重试循环

                except Exception as _retry_err:
                    if _attempt < max_retries:
                        self.log.warning(
                            f"Dify 调用失败 (第{_attempt}次)，正在重试: {_retry_err}"
                        )
                        await event_emitter(
                            {
                                "type": "status",
                                "data": {
                                    "description": f"调用出错，正在第{_attempt + 1}次重试...",
                                    "done": False,
                                },
                            }
                        )
                        await asyncio.sleep(3)
                    else:
                        raise

            # 2. 处理结果
            await event_emitter(
                {
                    "type": "status",
                    "data": {"description": "正在下载资源并合成文件...", "done": False},
                }
            )

            ppt_list = final_outputs.get("ppt_urls", [])
            img_bytes_list = []
            img_urls_list = []

            result = final_outputs.get("result", "")
            if "当前用户太多" in result and len(ppt_list) == 0:
                await self._clear_ppt_progress_embed(event_emitter, chat_id, message_id)
                await event_emitter(
                    {
                        "type": "status",
                        "data": {
                            "description": "抱歉！当前用户太多，请稍后尝试。",
                            "done": True,
                        },
                    }
                )
            else:
                # 并发下载图片
                async with httpx.AsyncClient(timeout=60) as dl_client:
                    for item in ppt_list:
                        url = item.get("ppt_url")
                        if not url:
                            continue
                        try:
                            if "https://dify.focusmedia.tech" in url:
                                url = url.replace(
                                    "https://dify.focusmedia.tech", self.valves.BASE_URL
                                )
                            resp = await dl_client.get(url)
                            if resp.status_code == 200:
                                image_data = resp.content
                                mime_type = (
                                    resp.headers.get("Content-Type") or "image/png"
                                )
                                # 上传以便预览
                                local_url = await self._upload_image_with_status(
                                    image_data,
                                    mime_type,
                                    request,
                                    __user__,
                                    event_emitter,
                                )
                                img_bytes_list.append(image_data)
                                img_urls_list.append(
                                    {
                                        "title": item.get("ppt_title", "Slide"),
                                        "url": local_url,
                                    }
                                )
                        except Exception as e:
                            print(f"DL Error: {e}")

                # 3. 构建 Viewer / Markdown
                if len(img_urls_list) > 0:
                    md_content = "\n\n"
                else:
                    class_name = final_outputs.get("class_name", "")
                    md_content = class_name

                # 4. 合成文件(ppt和pdf同时合成)
                dl_url_pdf = ""
                dl_url_ppt = ""
                if img_bytes_list:
                    try:
                        merged_data_pdf = (
                            await self._images_to_pdf_bytes_keep_resolution(
                                img_bytes_list
                            )
                        )
                        fname_pdf, fmime_pdf = (
                            f"ppt_{uuid.uuid4().hex[:6]}.pdf",
                            "application/pdf",
                        )
                        merged_data_ppt = (
                            await self._images_to_editable_ppt_bytes_fill_16_9(
                                img_bytes_list,
                            )
                        )
                        fname_ppt, fmime_ppt = (
                            f"ppt_{uuid.uuid4().hex[:6]}.pptx",
                            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )

                        # 【修改】并发上传 PDF+PPTX 到 OBS
                        pdf_object_key = f"ppt/output/{fname_pdf}"
                        ppt_object_key = f"ppt/output/{fname_ppt}"
                        dl_url_pdf, dl_url_ppt = await asyncio.gather(
                            self._upload_to_obs_async(
                                merged_data_pdf, pdf_object_key, fmime_pdf
                            ),
                            self._upload_to_obs_async(
                                merged_data_ppt, ppt_object_key, fmime_ppt
                            ),
                        )
                    except Exception as e:
                        self.log.error(f"Merge Error: {e}")
                        md_content += f"\n(文件合成出错: {e})"

                # 拼接尾部
                current_dify_id = str(uuid.uuid4())
                files_str = ",".join(map(str, processed_ids_list))
                hidden_tag = (
                    f"\n\n[dify_id]: http://context/{current_dify_id}?files={files_str}"
                )
                md_content = md_content + hidden_tag

                # 5. 清除进度条 + 发送结果（参考 seedance.py：DB清空 → viewer事件）
                #    不发空 embeds 事件，避免 socket handler 的 extend 逻辑恢复旧进度条
                await self._clear_ppt_progress_embed(event_emitter, chat_id, message_id)

                await event_emitter(
                    {"type": "message", "data": {"content": md_content}}
                )

                # 发送 PPT 查看器 embed
                viewer_html = None
                if len(img_urls_list) > 0:
                    viewer_html = self._build_ppt_viewer_html(
                        img_urls_list,
                        dl_url_pdf,
                        dl_url_ppt,
                        base_url=self._get_public_base_url(request),
                    )
                    await event_emitter(
                        {"type": "embeds", "data": {"embeds": [viewer_html]}}
                    )

                await event_emitter(
                    {
                        "type": "status",
                        "data": {
                            "description": f"任务完成，耗时{elapsed_time}s",
                            "done": True,
                        },
                    }
                )

                # 最终安全写入：直接设置 DB embeds 为最终值，
                # 防止前端 chatCompletedHandler 的 updateChatById 覆盖
                if chat_id and message_id:
                    try:
                        from open_webui.models.chats import Chats as _Chats

                        final_embeds = [viewer_html] if viewer_html else []
                        await _Chats.upsert_message_to_chat_by_id_and_message_id(
                            chat_id,
                            message_id,
                            {"embeds": final_embeds},
                        )
                    except Exception:
                        pass

        except Exception as e:
            import traceback

            error_trace = traceback.format_exc()
            self.log.error(f"Background Task Error: {str(e)}")
            print(
                f"Background Task Error===============================================: {e}"
            )
            print(
                f"Background Task Error===============================================: {repr(e)}",
                flush=True,
            )
            print(f"Traceback Details:\n{error_trace}", flush=True)
            await self._clear_ppt_progress_embed(event_emitter, chat_id, message_id)
            await event_emitter(
                {
                    "type": "status",
                    "data": {"description": f"任务出错: {str(e)}", "done": True},
                }
            )
            await event_emitter(
                {"type": "message", "data": {"content": f"\n\n**生成失败**: {e}"}}
            )

    async def _run_chatflow_background_task(
        self,
        dify_url: str,
        workflow_id: str,
        conversation_id: str,
        processed_ids_list: list,
        query: str,
        files: Union[List[dict], dict],
        request: Request,
        event_emitter: Callable,
        __user__=None,
        chat_id: str = "",
        message_id: str = "",
        session_id: str = "",
    ):
        try:
            user = await Users.get_user_by_id(__user__["id"])
            start_ts = int(time.time() * 1000)
            await self._emit_ppt_progress_once(
                event_emitter,
                start_ts,
                chat_id=chat_id,
                message_id=message_id,
            )

            headers = {
                "Authorization": f"Bearer {self.valves.API_KEY}",
                "Content-Type": "application/json",
            }

            payload = {
                "inputs": {},
                "query": query,
                "user": user.email,
                "response_mode": "streaming",
                # "conversation_id": "",
                # "files": []
            }
            if conversation_id:
                payload["conversation_id"] = conversation_id
            # 执行指定版本 Workflow
            if workflow_id:
                payload["workflow_id"] = workflow_id
            if (
                files
            ):  # 在 dify的 ChatFlow 应用中启用文件上传功能（确保聊天界面中可见附件图标）。  这个一定要启用，否则传不了文件。https://github.com/langgenius/dify/issues/24683
                payload["files"] = files

            print(f"payload:{payload}")

            final_ppt_outputs = {}
            full_answer_text = ""  # 用于累积 LLM 的回复文本
            has_error = False
            elapsed_time = 0

            current_dify_id = ""

            # 2. 发起流式请求，失败最多重试3次
            # [FIX] 显式拆分 timeout：read_timeout 设大以防 PPT 生成期间长时间无 SSE 事件导致连接被判定超时
            stream_timeout = httpx.Timeout(connect=30, read=1200, write=30, pool=30)
            max_retries = 3
            for _attempt in range(1, max_retries + 1):
                final_ppt_outputs = {}
                full_answer_text = ""
                has_error = False
                try:
                    async with httpx.AsyncClient(
                        timeout=stream_timeout, http2=False
                    ) as client:
                        async with client.stream(
                            "POST", dify_url, headers=headers, json=payload
                        ) as response:
                            if response.status_code != 200:
                                error_text = await response.aread()
                                raise Exception(
                                    f"Chat API Error {response.status_code}: {error_text}"
                                )

                            async for line in response.aiter_lines():
                                if not line or not line.startswith("data:"):
                                    continue

                                try:
                                    # 解析 SSE
                                    # [FIX] 解析 SSE - 兼容 "data: " 和 "data:" 两种格式
                                    json_str = line[5:].lstrip()
                                    data = json.loads(json_str)

                                    if not current_dify_id and data.get(
                                        "conversation_id"
                                    ):
                                        current_dify_id = data.get("conversation_id")

                                    event = data.get("event")

                                    # A. 处理文本流 (实时打字机效果)
                                    if event == "message":
                                        chunk_text = data.get("answer", "")
                                        full_answer_text += chunk_text

                                    # B. 处理节点状态 (用于显示进度)
                                    elif event == "node_started":
                                        node_data = data.get("data", {})
                                        title = node_data.get("title", "Processing")
                                        await event_emitter(
                                            {
                                                "type": "status",
                                                "data": {
                                                    "description": f"正在执行: {title}...",
                                                    "done": False,
                                                },
                                            }
                                        )

                                    # C. 处理工作流结束 (捕获 PPT 输出)
                                    elif event == "workflow_finished":
                                        workflow_data = data.get("data", {})
                                        if workflow_data.get("status") == "succeeded":
                                            outputs = workflow_data.get("outputs", {})
                                            elapsed_time = int(
                                                workflow_data.get("elapsed_time", 0)
                                            )
                                            if outputs:
                                                final_ppt_outputs = outputs
                                        else:
                                            print(
                                                f"Workflow finished with status: {workflow_data.get('status')}"
                                            )

                                    # D. 处理错误
                                    elif event == "error":
                                        has_error = True
                                        error_msg = data.get("message", "Unknown Error")
                                        raise Exception(f"Dify Error: {error_msg}")

                                except json.JSONDecodeError:
                                    print(
                                        f"[FIX] JSONDecodeError, raw line: {line[:500]}"
                                    )
                                except Exception:
                                    raise

                    if has_error:
                        raise Exception("Dify 流式响应包含错误")
                    break  # 成功，跳出重试循环

                except Exception as _retry_err:
                    if _attempt < max_retries:
                        self.log.warning(
                            f"Dify 调用失败 (第{_attempt}次)，正在重试: {_retry_err}"
                        )
                        await event_emitter(
                            {
                                "type": "status",
                                "data": {
                                    "description": f"调用出错，正在第{_attempt + 1}次重试...",
                                    "done": False,
                                },
                            }
                        )
                        await asyncio.sleep(3)
                    else:
                        # 最终失败，发送错误消息给用户
                        has_error = True
                        await event_emitter(
                            {
                                "type": "status",
                                "data": {
                                    "description": f"Error: {_retry_err}",
                                    "done": True,
                                },
                            }
                        )
                        await event_emitter(
                            {
                                "type": "message",
                                "data": {
                                    "content": f"\n\n **API Error**: {_retry_err}"
                                },
                            }
                        )

            # 3. 流结束后，检查是否有 PPT 数据需要处理

            final_ppt_outputs_answer = final_ppt_outputs.get("answer", "")
            final_ppt_urls = json.loads(full_answer_text.split("ppt_urls:")[-1].strip())

            ppt_images_bytes = []
            ppt_info_list = []
            if final_ppt_urls:
                await event_emitter(
                    {
                        "type": "status",
                        "data": {
                            "description": "PPT生成完毕，正在处理文件...",
                            "done": False,
                        },
                    }
                )
                # 下载图片
                async with httpx.AsyncClient(timeout=60) as dl_client:
                    for item in final_ppt_urls:
                        url = item.get("ppt_url")
                        if not url:
                            continue
                        try:
                            if (
                                "https://dify.focusmedia.tech" in url
                            ):  # 域名无法访问，ip可以
                                url = url.replace(
                                    "https://dify.focusmedia.tech", self.valves.BASE_URL
                                )
                            resp = await dl_client.get(url)
                            if resp.status_code == 200:
                                image_data = resp.content
                                mime_type = (
                                    resp.headers.get("Content-Type") or "image/png"
                                )
                                # 上传预览图
                                local_url = await self._upload_image_with_status(
                                    image_data,
                                    mime_type,
                                    request,
                                    __user__,
                                    event_emitter,
                                )
                                ppt_images_bytes.append(image_data)
                                ppt_info_list.append(
                                    {
                                        "title": item.get("ppt_title", "Slide"),
                                        "url": local_url,
                                    }
                                )
                        except:
                            pass

            # 构建 Viewer / Markdown
            md_content = "\n\n"

            # 合成文件
            dl_url_pdf = ""
            dl_url_ppt = ""
            if ppt_images_bytes:
                try:
                    merged_data_pdf = await self._images_to_pdf_bytes_keep_resolution(
                        ppt_images_bytes
                    )
                    fname_pdf, fmime_pdf = (
                        f"ppt_{uuid.uuid4().hex[:6]}.pdf",
                        "application/pdf",
                    )
                    merged_data_ppt = (
                        await self._images_to_editable_ppt_bytes_fill_16_9(
                            ppt_images_bytes,
                        )
                    )
                    fname_ppt, fmime_ppt = (
                        f"ppt_{uuid.uuid4().hex[:6]}.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )

                    # 【修改】并发上传 PDF+PPTX 到 OBS
                    pdf_object_key = f"ppt/output/{fname_pdf}"
                    ppt_object_key = f"ppt/output/{fname_ppt}"
                    dl_url_pdf, dl_url_ppt = await asyncio.gather(
                        self._upload_to_obs_async(
                            merged_data_pdf, pdf_object_key, fmime_pdf
                        ),
                        self._upload_to_obs_async(
                            merged_data_ppt, ppt_object_key, fmime_ppt
                        ),
                    )
                except Exception as e:
                    self.log.error(f"Merge Error: {e}")
                    md_content += f"\n(文件合成出错: {e})"

            # 拼接尾部
            files_str = ",".join(map(str, processed_ids_list))
            hidden_tag = (
                f"\n\n[dify_id]: http://context/{current_dify_id}?files={files_str}"
            )
            md_content = md_content + hidden_tag

            # 4. 清除进度条 + 发送结果（参考 seedance.py：DB清空 → viewer事件）
            await self._clear_ppt_progress_embed(event_emitter, chat_id, message_id)
            if not has_error:
                await event_emitter(
                    {"type": "message", "data": {"content": md_content}}
                )

                # 发送 PPT 查看器 embed
                viewer_html = None
                if len(ppt_info_list) > 0:
                    viewer_html = self._build_ppt_viewer_html(
                        ppt_info_list,
                        dl_url_pdf,
                        dl_url_ppt,
                        base_url=self._get_public_base_url(request),
                    )
                    await event_emitter(
                        {"type": "embeds", "data": {"embeds": [viewer_html]}}
                    )

                await event_emitter(
                    {
                        "type": "status",
                        "data": {
                            "description": f"生成完成，耗时{elapsed_time}s",
                            "done": True,
                        },
                    }
                )

                # 最终安全写入：直接设置 DB embeds 为最终值
                if chat_id and message_id:
                    try:
                        from open_webui.models.chats import Chats as _Chats

                        final_embeds = [viewer_html] if viewer_html else []
                        await _Chats.upsert_message_to_chat_by_id_and_message_id(
                            chat_id,
                            message_id,
                            {"embeds": final_embeds},
                        )
                    except Exception:
                        pass

        except Exception as e:
            self.log.error(f"Chatflow Background Error: {str(e)}")
            await self._clear_ppt_progress_embed(event_emitter, chat_id, message_id)
            await event_emitter(
                {
                    "type": "status",
                    "data": {"description": f"任务出错: {str(e)}", "done": True},
                }
            )
            await event_emitter(
                {
                    "type": "message",
                    "data": {"content": f"\n\n**System Error**: {str(e)}"},
                }
            )

    async def pipe(
        self,
        body: Dict[str, Any],
        __metadata__: dict[str, Any],
        __event_emitter__: Callable,
        __tools__: dict[str, Any] | None,
        __request__: Optional[Request] = None,
        __user__: Optional[dict] = None,
    ) -> Union[str, AsyncIterator[str]]:
        """
        Main method for sending requests to the Google Gemini endpoint.

        Args:
            body: The request body containing messages and other parameters.
            __metadata__: Request metadata
            __event_emitter__: Event emitter for status updates
            __tools__: Available tools
            __request__: FastAPI request object (for image upload)
            __user__: User information (for image upload)

        Returns:
            Response from Google Gemini API, which could be a string or an iterator for streaming.
        """
        # Setup logging for this request
        request_id = id(body)
        self.log.debug(f"Processing request {request_id}")
        self.log.debug(f"User request body: {__user__}")
        self.user = await Users.get_user_by_id(__user__["id"])

        # 自动为用户开启 iframeSandboxAllowSameOrigin，使 embed 可移至消息底部
        try:
            if self.user:
                _ui = (self.user.settings.ui if self.user.settings else None) or {}
                if not _ui.get("iframeSandboxAllowSameOrigin"):
                    _ui["iframeSandboxAllowSameOrigin"] = True
                    await Users.update_user_settings_by_id(__user__["id"], {"ui": _ui})
        except Exception:
            pass

        last_assistant_message = body["messages"][-1]

        # 获取用户消息中的图片和文本消息
        # 处理图片消息
        messages = body.get("messages", [])
        last_user_msg = next(
            (m for m in reversed(messages) if m.get("role") == "user"), None
        )

        if not last_user_msg:
            raise ValueError("No user message found")
        optimization_stats: List[Dict[str, Any]] = []
        prompt, current_images = await self._extract_images_from_message(
            last_user_msg, stats_list=optimization_stats
        )

        parent_message = __metadata__.get("parent_message", {})

        if "content" in parent_message:
            prompt = parent_message.get("content", "")

        try:
            dify_conversation_id, processed_file_ids = (
                self._extract_dify_id_from_history(body.get("messages", []))
            )
            self.log.debug(
                f"dify_conversation_id, processed_file_ids===================================================================={dify_conversation_id},{processed_file_ids}"
            )

            # 获取上传文档内容
            parent_message_files = parent_message.get("files", [])
            parent_message_files_id = (
                [pmf.get("id", "") for pmf in parent_message_files]
                if len(parent_message_files) > 0
                else []
            )

            files_to_send = []
            new_file_ids_record = set()  # 本次新增的 ID

            upload_document_contents = ""
            file_id_content_map = {}
            dify_upload_inputs = []
            files = __metadata__.get("files", [])
            if files:
                if len(files) > 0:
                    # 获取当前用户的headers和cookies
                    user_headers = {}
                    user_cookies = {}

                    user_cookies = __request__.cookies
                    auth = __request__.headers.get("authorization")
                    if auth:
                        user_headers["Authorization"] = auth
                    for file in files:
                        fid = file["id"]
                        if fid in parent_message_files_id:  # 代表当前对话需要上传的文件
                            # if fid not in processed_file_ids:
                            files_to_send.append(fid)
                            new_file_ids_record.add(str(fid))

                            file_url = file["url"]
                            content_type = file.get("content_type")
                            if file_url:
                                final_file_url = f"{self.valves.OPEN_WEBUI_BASE_URL}/api/v1/files/{fid}/content"

                                async def get_file_bytes_response():
                                    async with httpx.AsyncClient(timeout=60) as client:
                                        response = await client.get(
                                            final_file_url,
                                            headers=user_headers,
                                            cookies=user_cookies,
                                        )
                                        response.raise_for_status()
                                        # 获取 MIME 类型，如果没有则尝试从 URL 猜测
                                        mime_type = response.headers.get("Content-Type")
                                        if not mime_type:
                                            mime_type, _ = mimetypes.guess_type(
                                                final_file_url
                                            )

                                        # 返回二进制数据和类型，而不是 json
                                        return {
                                            "data": response.content,
                                            "mime_type": mime_type
                                            or content_type,  # "application/octet-stream"
                                        }

                                try:

                                    file_result = await self._retry_with_backoff(
                                        get_file_bytes_response
                                    )

                                    if file_result:
                                        # 将下载好的数据加入列表，准备传给 Dify
                                        file_id_content_map[fid] = fid
                                        dify_upload_inputs.append(
                                            {"inline_data": file_result}
                                        )
                                except Exception as e:
                                    # 下载失败记录日志，但不中断循环，继续处理下一个文件
                                    print(
                                        f"[Dify Tool] Failed to download file {fid}: {e}"
                                    )
            self.log.debug(
                f"files_to_send==============================================================={files_to_send}"
            )

            self.log.debug(
                f"new_file_ids_record==============================================================={new_file_ids_record}"
            )

            # 上传图片,文档到dify并获取id
            MAX_UPLOAD_LIMIT = 10
            n_files = len(dify_upload_inputs)
            n_imgs = len(current_images)
            total_count = n_files + n_imgs

            if total_count > MAX_UPLOAD_LIMIT:
                print(
                    f"[Dify] Upload items ({total_count}) exceed limit {MAX_UPLOAD_LIMIT}, truncating..."
                )

                half_limit = MAX_UPLOAD_LIMIT // 2  # 5

                # 策略：保护较少的一方，截断较多的一方

                # 情况1: 文件很少 (<=5)，保留所有文件，剩余名额给图片
                if n_files <= half_limit:
                    take_files = n_files
                    take_images = MAX_UPLOAD_LIMIT - n_files

                # 情况2: 图片很少 (<=5)，保留所有图片，剩余名额给文件
                elif n_imgs <= half_limit:
                    take_images = n_imgs
                    take_files = MAX_UPLOAD_LIMIT - n_imgs

                # 情况3: 两者都很多 (>5)，各取一半 (5 + 5)
                else:
                    take_files = half_limit
                    take_images = half_limit

                # 执行列表截取
                dify_upload_inputs = dify_upload_inputs[:take_files]
                current_images = current_images[:take_images]

                print(
                    f"[Dify] Truncated result: {len(dify_upload_inputs)} files, {len(current_images)} images."
                )

            # === 合并列表准备上传 ===
            # 将截断后的两个列表合并，统一传给 _upload_to_dify_direct
            all_items_to_upload = dify_upload_inputs + current_images

            user = await Users.get_user_by_id(__user__["id"])
            dify_files_list = None
            if len(all_items_to_upload) > 0:
                # current_images = current_images[:10]
                upload_result = await self._upload_to_dify_direct(
                    image_data=all_items_to_upload,
                    user_id=user.email,
                    api_base=self.valves.BASE_URL,
                    api_key=self.valves.API_KEY,
                    # 如果 image_data 是列表，mime_type 会自动从列表里取，这里不传也没关系
                    # 如果 image_data 是单张 raw bytes，建议传 mime_type
                    mime_type="image/png",
                )

                if isinstance(upload_result, list):
                    dify_files_list = upload_result
                elif isinstance(upload_result, dict):
                    dify_files_list = [upload_result]

            # 新逻辑
            final_query = prompt

            # if upload_document_contents:
            #     query = upload_document_contents + query
            # 我们要把 "旧的已处理ID" + "本次新增ID" 合并传给后台，让它写在最新的隐藏标签里
            updated_processed_ids = processed_file_ids.union(new_file_ids_record)

            # Non-streaming path (now also used for image generation)
            try:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "action": "image_processing",
                            "description": "正在处理PPT生成请求...",
                            "done": False,
                        },
                    }
                )

                # 3. 【关键修改】启动后台任务，并立即返回！
                # 这样 HTTP 请求会在 1秒内完成，Nginx 就不会报 503 了。
                chat_id = __metadata__.get("chat_id", "")
                message_id = __metadata__.get("message_id", "")
                session_id = __metadata__.get("session_id", "")
                if self.valves.APP_TYPE == "workflow":
                    dify_url = f"{self.valves.BASE_URL}/v1/workflows/run"
                    asyncio.create_task(
                        self._run_background_task(
                            dify_url=dify_url,
                            workflow_id=self.valves.WORKFLOW_ID,
                            processed_ids_list=list(updated_processed_ids),
                            query=final_query,
                            files=dify_files_list,
                            request=__request__,
                            event_emitter=__event_emitter__,
                            __user__=__user__,
                            chat_id=chat_id,
                            message_id=message_id,
                            session_id=session_id,
                        )
                    )
                else:
                    dify_url = f"{self.valves.BASE_URL}/v1/chat-messages"
                    asyncio.create_task(
                        self._run_chatflow_background_task(
                            dify_url=dify_url,
                            workflow_id=self.valves.WORKFLOW_ID,
                            conversation_id=dify_conversation_id,
                            processed_ids_list=list(updated_processed_ids),
                            query=final_query,
                            files=dify_files_list,
                            request=__request__,
                            event_emitter=__event_emitter__,
                            __user__=__user__,
                            chat_id=chat_id,
                            message_id=message_id,
                            session_id=session_id,
                        )
                    )

                # 4. 立即返回消息，告诉用户正在处理
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": "任务已提交，正在处理中...",
                            "done": False,
                        },
                    }
                )

            except Exception as e:
                self.log.exception(f"Error in non-streaming request {request_id}: {e}")
                return f"Error generating content: {e}"
                # yield f"Error generating content: {e}"

        except (ClientError, ServerError, APIError) as api_error:
            error_type = type(api_error).__name__
            error_msg = f"{error_type}: {api_error}"
            self.log.error(error_msg)
            return error_msg
            # yield error_msg

        except ValueError as ve:
            error_msg = f"Configuration error: {ve}"
            self.log.error(error_msg)
            return error_msg
            # yield error_msg

        except Exception as e:
            # Log the full error with traceback
            import traceback

            error_trace = traceback.format_exc()
            self.log.exception(f"Unexpected error: {e}\n{error_trace}")

            # Return a user-friendly error message
            return f"An error occurred while processing your request: {e}"
            # yield f"An error occurred while processing your request: {e}"
